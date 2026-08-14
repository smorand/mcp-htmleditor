"""Tests for the HTML to PPTX export.

The exporter is validated on three axes: slide detection (the historical bug),
content fidelity (tables, gantt, diagrams, images) and diagnostics (warnings
instead of silent success).
"""

from __future__ import annotations

import base64
import struct
import zlib
from pathlib import Path

import pytest
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches

from mcp_htmleditor.export.pptx_components import TableGrid, gantt_geometry
from mcp_htmleditor.export.pptx_style import (
    Box,
    StyleResolver,
    parse_color,
    parse_length,
    parse_pct,
    parse_px,
)
from mcp_htmleditor.export.to_pptx import (
    SLIDE_H_IN,
    SLIDE_W_IN,
    find_slides,
    to_pptx,
)

REFERENCE_DIR = Path(__file__).resolve().parents[1] / "templates" / "reference" / "slides"


# ---------------------------------------------------------------------------
# Test data
# ---------------------------------------------------------------------------


def _png(width: int, height: int, color: tuple[int, int, int] = (0, 58, 141)) -> bytes:
    """Build a minimal valid RGB PNG of the requested size."""

    def chunk(kind: bytes, payload: bytes) -> bytes:
        return (
            struct.pack(">I", len(payload))
            + kind
            + payload
            + struct.pack(">I", zlib.crc32(kind + payload) & 0xFFFFFFFF)
        )

    header = struct.pack(">2I5B", width, height, 8, 2, 0, 0, 0)
    raw = b"".join(b"\x00" + bytes(color) * width for _ in range(height))
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", header) + chunk(b"IDAT", zlib.compress(raw)) + chunk(b"IEND", b"")


def _data_uri(width: int = 8, height: int = 4) -> str:
    """Return a base64 PNG data URI usable as an ``<img src>``."""
    return "data:image/png;base64," + base64.b64encode(_png(width, height)).decode()


SECTION_DECK = """<!DOCTYPE html>
<html data-doc-type="presentation">
<head><meta charset="UTF-8"><style>.slide-h1 { color: #003A8D; }</style></head>
<body>
<section data-type="slide" data-id="s1" data-title="Intro">
  <h1 class="slide-h1">Title One</h1>
  <p>Some body text on the first slide.</p>
  <table>
    <thead><tr><th>Name</th><th>Value</th></tr></thead>
    <tbody>
      <tr><td>alpha</td><td>1</td></tr>
      <tr><td>beta</td><td>2</td></tr>
    </tbody>
  </table>
</section>
<section data-type="slide" data-id="s2" data-title="Charts">
  <div data-type="gantt">
    <div data-type="gantt-task" data-label="Design" data-start="2024-01" data-end="2024-02"
         data-color="#f00"></div>
    <div data-type="gantt-task" data-label="Build" data-start="2024-02" data-end="2024-04"></div>
  </div>
</section>
<section data-type="slide" data-id="s3" data-title="Diagram">
  <div data-type="arch-diagram" style="position:relative; height:200px;">
    <div data-type="arch-node" data-label="API" data-shape="box" data-x="10" data-y="10"
         data-width="20" data-height="15"
         style="left:10%; top:10%; border:2px solid #0f62fe; background:#edf5ff;">API</div>
    <div data-type="arch-node" data-label="DB" data-shape="cylinder" data-x="60" data-y="60">DB</div>
    <div class="arch-edge arch-line-h" data-type="arch-edge"
         style="left:30%; top:17%; width:30%; background:#003A8D;"></div>
    <div class="arch-edge arch-tip arch-tip-r"
         style="left:60%; top:17%; transform:translate(-100%,-50%);"></div>
    <div class="arch-edge arch-line-v" style="left:70%; top:25%; height:35%;"></div>
    <div class="arch-edge-label" style="left:45%; top:20%; transform:translateX(-50%);">HTTP</div>
    <div data-type="arch-edge" data-from="API" data-to="DB"
         style="left:35%; top:40%; width:10%;">&rarr;</div>
  </div>
</section>
<section data-type="slide" data-id="s4" data-title="Spans">
  <table data-type="table">
    <colgroup><col style="width:50%"><col style="width:25%"><col style="width:25%"></colgroup>
    <thead><tr><th>A</th><th>B</th><th>C</th></tr></thead>
    <tbody>
      <tr><td rowspan="2">merged down</td><td colspan="2">merged across</td></tr>
      <tr><td>b2</td><td>c2</td></tr>
    </tbody>
  </table>
</section>
</body>
</html>
"""


EI_DECK = """<!DOCTYPE html>
<html lang="fr" data-doc-type="presentation">
<head><meta charset="UTF-8"><style>
  :root {{ --ei-blue: #003A8D; --ei-orange: #FBAE40; }}
  .slide-inner {{ background: #ffffff; }}
  .slide-title-rule {{ height: 3px; width: 64px; background: var(--ei-orange); }}
  .gantt-bar {{ color: #ffffff; }}
</style></head>
<body>
<div class="toolbar"><span>Diapositive :</span><select><option>01</option></select></div>
<article class="slide" data-type="slide" data-slide-type="title" data-id="slide-0"
         data-title="Couverture">
  <img class="slide-cover-img" src="{cover}">
  <div class="slide-cover-body">
    <div class="slide-cover-title">Titre de la presentation</div>
    <div class="slide-cover-subtitle">Sous-titre</div>
  </div>
  <div class="slide-cover-logos">
    <div class="logos-left"><img class="logo-cm" src="{logo}"><img class="logo-cic" src="{logo}"></div>
    <img class="logo-ei" src="{logo}">
  </div>
</article>
<article class="slide" data-type="slide" data-slide-type="section" data-id="slide-1"
         data-title="Section">
  <div class="slide-section-band"></div>
  <div class="slide-section-body">
    <div class="slide-section-num">Section 01</div>
    <div class="slide-section-title">Contexte</div>
  </div>
</article>
<article class="slide" data-type="slide" data-slide-type="content" data-id="slide-2"
         data-title="Planning">
  <div class="slide-inner">
    <div class="slide-eyebrow">Planning &middot; Slide 03 / 03</div>
    <h1 class="slide-h1">Planning <span>2026</span></h1>
    <div class="slide-title-rule"></div>
    <div class="slide-body">
      <div data-type="gantt" style="width:100%;">
        <div class="gantt-head">
          <div class="gantt-label">Chantier</div>
          <div class="gantt-track"><div>T1</div><div>T2</div></div>
          <div class="gantt-dates">Periode</div>
        </div>
        <div class="gantt-row">
          <div class="gantt-label">Cadrage</div>
          <div class="gantt-track">
            <div class="gantt-bar" data-type="gantt-task" data-id="t1" data-label="Cadrage"
                 data-start="2026-01" data-end="2026-02" data-color="#003A8D"
                 style="left:0%; width:50%; background:#003A8D;">Jan &rarr; Fev</div>
          </div>
          <div class="gantt-dates">Jan a Fev</div>
        </div>
        <div class="gantt-legend">
          <span><i style="background:#003A8D;"></i>Socle</span>
        </div>
      </div>
      <div class="cds-grid cols-2">
        <div class="cds-tile">
          <div class="tile-eyebrow">Enjeu 01</div>
          <div class="tile-title">Productivite</div>
          <ul><li>un point</li><li>un autre</li></ul>
        </div>
        <div class="cds-tile"><div class="tile-title">Charte</div><p>texte</p></div>
      </div>
    </div>
  </div>
  <div class="slide-foot">
    <div class="slide-foot-logo"><span class="logo-disc"><img src="{logo}"></span></div>
    <span class="slide-foot-page">3</span>
    <span class="slide-foot-title">Comite de pilotage</span>
  </div>
</article>
<script>const TOTAL = 3; function render() {{ return TOTAL; }}</script>
</body>
</html>
"""


@pytest.fixture
def section_deck(tmp_path: Path) -> Path:
    """Write the legacy ``<section data-type="slide">`` deck to disk."""
    path = tmp_path / "deck.html"
    path.write_text(SECTION_DECK, encoding="utf-8")
    return path


@pytest.fixture
def ei_deck(tmp_path: Path) -> Path:
    """Write a compact Euro-Information deck exercising the whole charter."""
    path = tmp_path / "ei.html"
    path.write_text(EI_DECK.format(cover=_data_uri(16, 9), logo=_data_uri(4, 4)), encoding="utf-8")
    return path


def _all_text(prs: Presentation) -> str:
    """Concatenate every piece of text of a presentation, tables included."""
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.extend(cell.text for cell in row.cells)
    return "\n".join(chunks)


def _pictures(prs: Presentation) -> list[object]:
    """Return every picture shape of a presentation."""
    return [shape for slide in prs.slides for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]


# ---------------------------------------------------------------------------
# Slide detection
# ---------------------------------------------------------------------------


def test_find_slides_matches_any_tag_with_the_attribute() -> None:
    """Detection is attribute based, so <article> counts like <section>."""
    soup = BeautifulSoup(
        '<body><article data-type="slide">a</article><section data-type="slide">b</section></body>',
        "html.parser",
    )
    assert len(find_slides(soup)) == 2


def test_find_slides_falls_back_on_the_slide_class() -> None:
    """Templates without the attribute are detected through ``.slide``."""
    soup = BeautifulSoup(
        '<body><article class="slide active" id="slide-0">a</article>'
        '<article class="slide" id="slide-1">b</article></body>',
        "html.parser",
    )
    assert len(find_slides(soup)) == 2


def test_find_slides_ignores_nested_slides() -> None:
    """A slide nested in another slide is not counted twice."""
    soup = BeautifulSoup(
        '<article data-type="slide"><div data-type="slide">x</div></article>',
        "html.parser",
    )
    assert len(find_slides(soup)) == 1


def test_euro_information_template_exports_every_slide(tmp_path: Path) -> None:
    """The EI reference template exports all of its <article> slides."""
    source = REFERENCE_DIR / "euro-information.html"
    expected = len(find_slides(BeautifulSoup(source.read_text(encoding="utf-8"), "html.parser")))
    out = tmp_path / "ei.pptx"

    report = to_pptx(str(source), str(out))

    assert expected == 3
    assert report.slide_count == expected
    assert len(Presentation(str(out)).slides) == expected


def test_ibm_carbon_template_exports_nine_slides(tmp_path: Path) -> None:
    """The Carbon reference template has 9 slides and no data-type attribute."""
    out = tmp_path / "carbon.pptx"

    report = to_pptx(str(REFERENCE_DIR / "ibm-carbon.html"), str(out))

    assert report.slide_count == 9


def test_section_deck_does_not_regress(section_deck: Path, tmp_path: Path) -> None:
    """The legacy <section> markup still exports one slide per section."""
    out = tmp_path / "out" / "deck.pptx"

    report = to_pptx(str(section_deck), str(out))

    prs = Presentation(str(out))
    assert report.slide_count == 4
    assert len(prs.slides) == 4
    assert "Title One" in _all_text(prs)
    assert any(shape.has_table for shape in prs.slides[0].shapes)


def test_slides_are_sixteen_by_nine(section_deck: Path, tmp_path: Path) -> None:
    """Slides use the 13.333 x 7.5 inch canvas of the templates."""
    out = tmp_path / "deck.pptx"
    to_pptx(str(section_deck), str(out))

    prs = Presentation(str(out))
    assert prs.slide_width == Inches(SLIDE_W_IN)
    assert prs.slide_height == Inches(SLIDE_H_IN)


# ---------------------------------------------------------------------------
# Shell and code exclusion
# ---------------------------------------------------------------------------


def test_script_and_shell_are_never_exported(tmp_path: Path) -> None:
    """Navigation code and shell widgets stay out of the PPTX."""
    out = tmp_path / "carbon.pptx"
    to_pptx(str(REFERENCE_DIR / "ibm-carbon.html"), str(out))

    text = _all_text(Presentation(str(out)))
    for forbidden in ("const TOTAL", "function ", "addEventListener", "querySelector", "=>"):
        assert forbidden not in text
    assert "Diapositive" not in text  # toolbar label of the navigation shell


def test_document_without_slides_warns_and_keeps_content(tmp_path: Path) -> None:
    """A document with no slide element yields one slide plus a warning."""
    source = tmp_path / "plain.html"
    source.write_text(
        "<html><body><h1>Just a doc</h1><p>text</p><script>const x = 1;</script></body></html>",
        encoding="utf-8",
    )
    out = tmp_path / "plain.pptx"

    report = to_pptx(str(source), str(out))

    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    assert report.warnings
    assert "Just a doc" in _all_text(prs)
    assert "const x" not in _all_text(prs)


# ---------------------------------------------------------------------------
# Images
# ---------------------------------------------------------------------------


def test_base64_images_are_embedded(tmp_path: Path) -> None:
    """A base64 data URI ends up as a real picture in the deck."""
    source = tmp_path / "img.html"
    source.write_text(
        '<html><body><article data-type="slide" data-slide-type="content">'
        f'<div class="slide-body"><img src="{_data_uri()}"></div>'
        "</article></body></html>",
        encoding="utf-8",
    )
    out = tmp_path / "img.pptx"

    report = to_pptx(str(source), str(out))

    assert report.warnings == []
    assert len(_pictures(Presentation(str(out)))) == 1


def test_relative_images_resolve_against_the_html_file(tmp_path: Path, monkeypatch) -> None:
    """A relative src is read next to the HTML file, not next to the CWD."""
    assets = tmp_path / "deck"
    assets.mkdir()
    (assets / "chart.png").write_bytes(_png(8, 4))
    source = assets / "deck.html"
    source.write_text(
        '<html><body><article data-type="slide"><img src="chart.png"></article></body></html>',
        encoding="utf-8",
    )
    elsewhere = tmp_path / "elsewhere"
    elsewhere.mkdir()
    monkeypatch.chdir(elsewhere)
    out = tmp_path / "rel.pptx"

    report = to_pptx(str(source), str(out))

    assert report.warnings == []
    assert len(_pictures(Presentation(str(out)))) == 1


def test_missing_and_remote_images_are_reported(tmp_path: Path) -> None:
    """Unresolvable images are skipped with an explicit warning each."""
    source = tmp_path / "missing.html"
    source.write_text(
        '<html><body><article data-type="slide">'
        '<img src="nope.png"><img src="https://example.com/x.png">'
        "</article></body></html>",
        encoding="utf-8",
    )
    out = tmp_path / "missing.pptx"

    report = to_pptx(str(source), str(out))

    assert not _pictures(Presentation(str(out)))
    assert any("introuvable" in warning for warning in report.warnings)
    assert any("distante" in warning for warning in report.warnings)


def test_annotations_are_positioned_inside_the_image(tmp_path: Path) -> None:
    """Annotation percentages resolve in the image box, not the whole slide."""
    source = tmp_path / "annot.html"
    source.write_text(
        '<html><body><article data-type="slide" data-slide-type="content">'
        '<div class="slide-body">'
        '<div data-type="annotated-image" style="max-width:400px; margin:0 auto;">'
        f'<img src="{_data_uri(16, 9)}" style="width:100%">'
        '<div data-type="annotation" data-x="50" data-y="50" '
        'style="left:50%; top:50%; background:#003A8D; color:#fff">middle</div>'
        "</div></div></article></body></html>",
        encoding="utf-8",
    )
    out = tmp_path / "annot.pptx"

    to_pptx(str(source), str(out))

    slide = Presentation(str(out)).slides[0]
    picture = next(s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    annotation = next(s for s in slide.shapes if s.has_text_frame and s.text_frame.text == "middle")
    assert picture.left <= annotation.left <= picture.left + picture.width
    assert picture.top <= annotation.top <= picture.top + picture.height


# ---------------------------------------------------------------------------
# Components
# ---------------------------------------------------------------------------


def test_gantt_bars_are_shapes_not_a_table(section_deck: Path, tmp_path: Path) -> None:
    """Tasks become positioned bars; no fallback table is produced."""
    out = tmp_path / "deck.pptx"
    to_pptx(str(section_deck), str(out))

    slide = Presentation(str(out)).slides[1]
    assert not any(shape.has_table for shape in slide.shapes)
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    assert "Design" in texts
    assert "Build" in texts


def test_gantt_geometry_from_dates() -> None:
    """Without inline positioning, bars are placed from data-start/data-end."""
    soup = BeautifulSoup(
        '<div data-type="gantt-task" data-start="2024-04" data-end="2024-06"></div>',
        "html.parser",
    )
    task = soup.div
    assert task is not None
    left, width = gantt_geometry(task, (_month(2024, 1), 12))
    assert left == pytest.approx(25.0)
    assert width == pytest.approx(25.0)


def _month(year: int, month: int) -> int:
    """Absolute month index used by the Gantt period helpers."""
    return year * 12 + month - 1


def test_gantt_geometry_prefers_inline_percentages() -> None:
    """Inline left/width percentages win, they are what the browser renders."""
    soup = BeautifulSoup(
        '<div data-type="gantt-task" data-start="2024-01" data-end="2024-12" style="margin-left:10%; width:30%"></div>',
        "html.parser",
    )
    task = soup.div
    assert task is not None
    assert gantt_geometry(task, (_month(2024, 1), 12)) == (10.0, 30.0)


# A hand-authored Gantt with no data-type/gantt-row markup: two label+track rows,
# stacked sub-lane bars (top offsets), a milestone marker line, a hatch-patterned
# "done" bar, a header row and a legend row as *siblings* of the row stack (not
# descendants), matching the shape found in the client deck this variant fixes.
INLINE_GANTT_DECK = """<!DOCTYPE html>
<html data-doc-type="presentation">
<body>
<section data-type="slide" data-id="s1" data-title="Roadmap">
  <div class="slide-body">
    <div style="display:flex; margin-left:80px;">
      <div style="flex:0 0 33.33%; text-align:center;">Jan 26</div>
      <div style="flex:0 0 33.33%; text-align:center;">Feb 26</div>
      <div style="flex:0 0 33.33%; text-align:center;">Mar 26</div>
    </div>
    <div style="flex:1;">
      <div style="display:flex; align-items:stretch; min-height:26px;">
        <div style="flex:0 0 80px;">Alpha</div>
        <div style="flex:1; position:relative; min-height:26px;">
          <div style="position:absolute; left:33%; top:0; bottom:0; width:2px; background:#FBAE40;"></div>
          <div style="position:absolute; left:0%; width:30%; top:2px; height:10px; background:#003A8D;">Design</div>
          <div style="position:absolute; left:35%; width:30%; top:13px; height:10px; background:#888;
                      background-image:repeating-linear-gradient(45deg,rgba(255,255,255,.45) 0,
                      rgba(255,255,255,.45) 2px,transparent 2px,transparent 5px);">Done</div>
        </div>
      </div>
      <div style="display:flex; align-items:stretch; min-height:15px;">
        <div style="flex:0 0 80px;">Beta</div>
        <div style="flex:1; position:relative; min-height:15px;">
          <div style="position:absolute; left:60%; width:25%; top:2px; height:10px; background:#C0392B;">Build</div>
        </div>
      </div>
    </div>
    <div style="display:flex; flex-wrap:wrap; gap:5px;">
      <span><span style="width:10px; height:8px; background:#003A8D; display:inline-block;"></span>Alpha</span>
      <span><span style="width:10px; height:8px; background:#C0392B; display:inline-block;"></span>Beta</span>
      <span><span style="width:10px; height:8px; background:#888;
            background-image:repeating-linear-gradient(45deg,rgba(255,255,255,.45) 0,
            rgba(255,255,255,.45) 2px,transparent 2px,transparent 5px); display:inline-block;"></span>Termine</span>
      <span style="border-left:2px solid #d9dee6; padding-left:8px;">
        <span style="width:2px; height:11px; background:#FBAE40; display:inline-block;"></span>Aujourd'hui
      </span>
    </div>
  </div>
</section>
</body>
</html>
"""


@pytest.fixture
def inline_gantt_deck(tmp_path: Path) -> Path:
    """Write the hand-authored, class-free Gantt deck to disk."""
    path = tmp_path / "inline-gantt.html"
    path.write_text(INLINE_GANTT_DECK, encoding="utf-8")
    return path


def test_inline_gantt_bars_are_positioned_shapes(inline_gantt_deck: Path, tmp_path: Path) -> None:
    """A class-free Gantt is detected structurally: real bars, no fallback table,
    no giant flattened text block for the whole chart (the original bug)."""
    out = tmp_path / "deck.pptx"
    report = to_pptx(str(inline_gantt_deck), str(out))

    assert report.warnings == []
    slide = Presentation(str(out)).slides[0]
    assert not any(shape.has_table for shape in slide.shapes)
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    assert "Design" in texts
    assert "Build" in texts
    assert "Done" in texts


def test_inline_gantt_stacked_bars_do_not_overlap(inline_gantt_deck: Path, tmp_path: Path) -> None:
    """Two bars sharing a row (different inline top offsets) land at distinct
    vertical bands instead of both centering on the row's mid-line."""
    out = tmp_path / "deck.pptx"
    to_pptx(str(inline_gantt_deck), str(out))

    slide = Presentation(str(out)).slides[0]
    design = next(s for s in slide.shapes if s.has_text_frame and s.text_frame.text == "Design")
    done = next(s for s in slide.shapes if s.has_text_frame and s.text_frame.text == "Done")
    assert design.top != done.top
    assert design.top + design.height <= done.top + Inches(0.02)


def test_inline_gantt_marker_line_is_a_thin_colored_rect(inline_gantt_deck: Path, tmp_path: Path) -> None:
    """A markerless absolutely positioned track child (no text) becomes a thin
    vertical rect at its left percentage, reusing the grid-line primitive."""
    out = tmp_path / "deck.pptx"
    to_pptx(str(inline_gantt_deck), str(out))

    slide = Presentation(str(out)).slides[0]
    markers = [
        s
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and s.fill.type is not None
        and str(s.fill.fore_color.rgb) == "FBAE40"
        and s.width < Inches(0.05)
    ]
    # one is the marker line drawn across row 1's track, the other the legend swatch
    assert len(markers) == 2
    tops = sorted(s.top for s in markers)
    assert tops[1] - tops[0] > Inches(0.5)  # the track marker sits well above the legend row


def test_inline_gantt_hatch_marks_the_done_bar(inline_gantt_deck: Path, tmp_path: Path) -> None:
    """A CSS repeating-linear-gradient background becomes a real PPTX pattern
    fill, both on the bar and on the matching legend swatch."""
    out = tmp_path / "deck.pptx"
    to_pptx(str(inline_gantt_deck), str(out))

    from pptx.enum.dml import MSO_FILL_TYPE

    slide = Presentation(str(out)).slides[0]
    hatched_shapes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.fill.type == MSO_FILL_TYPE.PATTERNED
    ]
    assert len(hatched_shapes) == 2  # the "Done" bar and the "Termine" legend swatch
    done_label = next(s for s in slide.shapes if s.has_text_frame and s.text_frame.text == "Done")
    bar = next(
        s
        for s in hatched_shapes
        if s.top <= done_label.top <= s.top + s.height and s.left <= done_label.left <= s.left + s.width
    )
    assert bar is not None
    assert "Termine" in [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]


def test_inline_gantt_header_and_legend_render_once(inline_gantt_deck: Path, tmp_path: Path) -> None:
    """The header (month cells) and legend, siblings of the row stack in the
    DOM, are consumed before the generic block-flow walk reaches them: no
    duplicate flattened text block, no leftover generic paragraph."""
    out = tmp_path / "deck.pptx"
    to_pptx(str(inline_gantt_deck), str(out))

    slide = Presentation(str(out)).slides[0]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    assert texts.count("Jan 26") == 1
    assert not any("Jan 26" in t and "Feb 26" in t and "\n" in t for t in texts)
    assert "Aujourd'hui" in "".join(texts)
    assert "Alpha" in texts and texts.count("Alpha") == 2  # row label and legend entry


def test_inline_gantt_long_label_is_truncated_not_overflowing(tmp_path: Path) -> None:
    """A label far longer than its narrow bar is shrunk then ellipsised, instead
    of word-wrapping past the bar's height into the row above/below it (the
    ghost-text bug the fixed-width, non-wrapping label textbox now avoids)."""
    deck = INLINE_GANTT_DECK.replace(
        'left:0%; width:30%; top:2px; height:10px; background:#003A8D;">Design',
        'left:0%; width:6%; top:2px; height:10px; background:#003A8D;">'
        "A very long task label that cannot possibly fit in this narrow bar",
    )
    source = tmp_path / "long-label.html"
    source.write_text(deck, encoding="utf-8")
    out = tmp_path / "deck.pptx"

    to_pptx(str(source), str(out))

    slide = Presentation(str(out)).slides[0]
    label = next(s for s in slide.shapes if s.has_text_frame and s.text_frame.text.startswith("A very long"))
    assert label.text_frame.word_wrap is False
    assert label.text_frame.text.endswith("\u2026")
    assert len(label.text_frame.text) < len("A very long task label that cannot possibly fit in this narrow bar")


def test_arch_nodes_become_autoshapes(section_deck: Path, tmp_path: Path) -> None:
    """Diagram nodes are real shapes with an outline, and keep their label."""
    out = tmp_path / "deck.pptx"
    to_pptx(str(section_deck), str(out))

    slide = Presentation(str(out)).slides[2]
    shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    labels = {s.text_frame.text for s in shapes if s.has_text_frame}
    assert {"API", "DB"} <= labels
    node = next(s for s in shapes if s.has_text_frame and s.text_frame.text == "API")
    assert str(node.line.color.rgb) == "0F62FE"


def test_arch_edges_become_segments_tips_and_labels(section_deck: Path, tmp_path: Path) -> None:
    """CSS connectors become thin rectangles, arrow heads and labels."""
    out = tmp_path / "deck.pptx"
    to_pptx(str(section_deck), str(out))

    slide = Presentation(str(out)).slides[2]
    shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    segments = [s for s in shapes if s.auto_shape_type == MSO_SHAPE.RECTANGLE]
    tips = [s for s in shapes if s.auto_shape_type == MSO_SHAPE.ISOSCELES_TRIANGLE]
    assert len(segments) >= 2  # one horizontal, one vertical
    assert len(tips) == 1
    assert tips[0].rotation == 90
    texts = {s.text_frame.text for s in slide.shapes if s.has_text_frame}
    assert "HTTP" in texts
    assert "\u2192" in texts


def test_table_merges_spans(section_deck: Path, tmp_path: Path) -> None:
    """colspan and rowspan are converted to real PPTX merges."""
    out = tmp_path / "deck.pptx"
    to_pptx(str(section_deck), str(out))

    table = next(s.table for s in Presentation(str(out)).slides[3].shapes if s.has_table)
    assert len(table.rows) == 3
    assert len(table.columns) == 3
    assert table.cell(1, 0).is_merge_origin
    assert table.cell(2, 0).is_spanned
    assert table.cell(1, 1).is_merge_origin
    assert table.cell(1, 2).is_spanned
    assert "merged across" in table.cell(1, 1).text


def test_table_grid_resolves_spans() -> None:
    """The occupancy grid keeps every cell reachable at its own origin."""
    soup = BeautifulSoup(
        "<table><tr><th>a</th><th>b</th><th>c</th></tr>"
        '<tr><td rowspan="2">x</td><td colspan="2">y</td></tr>'
        "<tr><td>p</td><td>q</td></tr></table>",
        "html.parser",
    )
    table = soup.table
    assert table is not None
    grid = TableGrid(table)
    assert grid.column_count == 3
    assert grid.header_rows == {0}
    assert grid.rows[1][1] is grid.rows[1][2]
    assert grid.rows[2][0] is grid.rows[1][0]


def test_table_grid_uses_colgroup_widths() -> None:
    """Declared colgroup widths drive the PPTX column widths."""
    soup = BeautifulSoup(
        '<table><colgroup><col style="width:50%"><col style="width:25%">'
        '<col style="width:25%"></colgroup>'
        "<tr><td>a</td><td>b</td><td>c</td></tr></table>",
        "html.parser",
    )
    table = soup.table
    assert table is not None
    widths = TableGrid(table).column_widths(10.0)
    assert widths == pytest.approx([5.0, 2.5, 2.5])


def test_table_grid_keeps_numbering_columns_narrow() -> None:
    """An agenda numbering column keeps its template width."""
    soup = BeautifulSoup(
        '<table><tr><td class="num">01</td><td>label</td></tr></table>',
        "html.parser",
    )
    table = soup.table
    assert table is not None
    widths = TableGrid(table).column_widths(10.0)
    assert widths[0] < 1.0
    assert widths[1] > 8.0


# ---------------------------------------------------------------------------
# Geometry and style helpers
# ---------------------------------------------------------------------------


def test_parse_lengths() -> None:
    """Percentages and pixels parse, anything else is 0 or None."""
    assert parse_pct("25%") == 25.0
    assert parse_pct("nope") == 0.0
    assert parse_px(" 48.5px ") == 48.5
    assert parse_px("auto") == 0.0
    assert parse_length("50%", 10.0) == pytest.approx(5.0)
    assert parse_length("72px", 10.0) == pytest.approx(1.0, abs=1e-3)
    assert parse_length("auto", 10.0) is None
    assert parse_length(None, 10.0) is None


def test_box_arithmetic() -> None:
    """Percentage sub-boxes and insets stay inside the parent box."""
    box = Box(1.0, 2.0, 10.0, 4.0)
    inner = box.pct(50.0, 50.0, 50.0, 50.0)
    assert (inner.left, inner.top, inner.width, inner.height) == (6.0, 4.0, 5.0, 2.0)
    inset = box.inset(0.5)
    assert (inset.left, inset.width) == (1.5, 9.0)
    assert box.center_horizontally(2.0).left == pytest.approx(5.0)
    assert box.right == 11.0
    assert box.bottom == 6.0


def test_parse_color_forms() -> None:
    """Hex, short hex, rgb(), named colors and var() all resolve."""
    assert parse_color("#003A8D") == "003A8D"
    assert parse_color("#abc") == "AABBCC"
    assert parse_color("rgba(255,255,0,.92)") == "FFFF00"
    assert parse_color("white") == "FFFFFF"
    assert parse_color("var(--ei-blue)", {"--ei-blue": "#003A8D"}) == "003A8D"
    assert parse_color("linear-gradient(135deg, #1A1A2E 0%, #16213E 60%)") == "1A1A2E"
    assert parse_color(None) is None


def test_style_resolver_indexes_classes_without_leaking_compounds() -> None:
    """``.notif.error`` must not paint every ``.notif``, comments are ignored."""
    soup = BeautifulSoup(
        "<html><head><style>"
        "/* a comment with a * star */\n"
        ".notif { background: #EDF5FF; }\n"
        ".notif.error { background: #FFF1F1; }\n"
        ".card .name { color: #003A8D; font-weight: 700; }\n"
        ".slide-header { border-bottom: 4px solid #0F62FE; }\n"
        "</style></head>"
        '<body><div class="notif error">x</div></body></html>',
        "html.parser",
    )
    resolver = StyleResolver.from_soup(soup)

    assert resolver.class_props["notif"]["background"] == "#EDF5FF"
    assert "error" not in resolver.class_props
    assert resolver.class_props["name"]["color"] == "#003A8D"
    node = soup.find("div")
    assert node is not None
    assert resolver.color(node, ("background",)) == "EDF5FF"


def test_style_resolver_applies_inline_overrides() -> None:
    """Inline declarations win over the class and tag scales."""
    soup = BeautifulSoup(
        '<p class="slide-subtitle" style="font-size:20px; color:#FF0000; font-weight:700; text-align:center">x</p>',
        "html.parser",
    )
    node = soup.find("p")
    assert node is not None
    style = StyleResolver.from_soup(soup).style(node)
    assert style.size == 20.0
    assert style.color == "FF0000"
    assert style.bold is True
    assert style.align == "center"


# ---------------------------------------------------------------------------
# Euro-Information charter
# ---------------------------------------------------------------------------


def test_ei_charter_is_drawn(ei_deck: Path, tmp_path: Path) -> None:
    """Every EI slide type gets its frame, its band and its footer ring."""
    out = tmp_path / "ei.pptx"
    report = to_pptx(str(ei_deck), str(out))

    prs = Presentation(str(out))
    assert report.slide_count == 3
    assert report.warnings == []

    fills = {
        str(shape.fill.fore_color.rgb)
        for slide in prs.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.fill.type is not None
    }
    assert "003A8D" in fills  # slide frame and section background
    assert "284AAA" in fills  # section band
    assert "FBAE40" in fills  # title accent rule

    ovals = [
        shape
        for shape in prs.slides[2].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.auto_shape_type == MSO_SHAPE.OVAL
    ]
    assert len(ovals) == 1  # the footer logo ring
    # the footer title is uppercased by the charter (text-transform: uppercase)
    assert "COMITE DE PILOTAGE" in _all_text(prs)


def test_ei_cover_places_every_logo(ei_deck: Path, tmp_path: Path) -> None:
    """Cover pictures are the image plus the three logos, none overlapping."""
    out = tmp_path / "ei.pptx"
    to_pptx(str(ei_deck), str(out))

    slide = Presentation(str(out)).slides[0]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 4
    left_logos = sorted((s.left, s.width) for s in pictures if s.top > Inches(6.0) and s.left < Inches(6.0))
    assert len(left_logos) == 2
    first, second = left_logos
    assert first[0] + first[1] <= second[0]


def test_ei_gantt_keeps_scale_labels_and_legend(ei_deck: Path, tmp_path: Path) -> None:
    """The Gantt scale, the dates column and the legend all reach the slide."""
    out = tmp_path / "ei.pptx"
    to_pptx(str(ei_deck), str(out))

    texts = {shape.text_frame.text for shape in Presentation(str(out)).slides[2].shapes if shape.has_text_frame}
    assert {"T1", "T2", "Cadrage", "Jan a Fev", "Socle"} <= texts
    assert any(text.startswith("Jan") and "\u2192" in text for text in texts)


def test_ei_tiles_keep_their_bullets(ei_deck: Path, tmp_path: Path) -> None:
    """List items inside a tile are exported with an explicit bullet marker."""
    out = tmp_path / "ei.pptx"
    to_pptx(str(ei_deck), str(out))

    text = _all_text(Presentation(str(out)))
    assert "\u2022 un point" in text
    assert "\u2022 un autre" in text
