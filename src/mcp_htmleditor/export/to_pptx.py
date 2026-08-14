"""Export an mcp-htmleditor HTML presentation to PPTX with python-pptx.

Design
------
Every element carrying ``data-type="slide"`` becomes one 16:9 slide
(13.333 x 7.5 inches), which is the reference canvas of the templates
(960 x 540 CSS pixels, so one CSS pixel is one point). Nothing outside the
slides is exported: the navigation shell, ``<script>`` and ``<style>`` blocks
are skipped, which is what kept the previous version from being usable.

Inside a slide the exporter runs a small block flow engine:

1. the template chrome is drawn first (backgrounds, frames, footers) from
   ``data-slide-type`` and from the charter detected on the document;
2. the slide subtree is walked recursively and turned into an ordered list of
   blocks (text, tile grid, notification, table, gantt, arch diagram,
   annotated image, image, accent rule);
3. blocks are stacked vertically in their region, flexible blocks absorb the
   remaining space, and everything is scaled down when the content overflows.

Fidelity notes live in ``skill/workflow-export.md``.
"""

from __future__ import annotations

import base64
import binascii
import io
import logging
import math
import re
from collections.abc import Callable
from dataclasses import dataclass, field, replace
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_PATTERN_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from ..tracing import trace_span
from .pptx_components import (
    GANTT_DATES_PX,
    GANTT_LABEL_PX,
    GanttRow,
    TableGrid,
    apply_transform,
    color_of,
    drop_theme_style,
    first_pct,
    gantt_geometry,
    gantt_period,
    gantt_task_band,
    inches,
    is_light,
    set_cell_border,
    split_runs,
)
from .pptx_style import (
    PX_IN,
    SLIDE_H_IN,
    SLIDE_W_IN,
    Box,
    StyleResolver,
    TextStyle,
    Theme,
    block_height,
    classes,
    has_class,
    parse_color,
    parse_length,
    parse_pct,
    parse_px,
    style_props,
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SKIP_TAGS = frozenset(
    {
        "script",
        "style",
        "noscript",
        "template",
        "link",
        "meta",
        "svg",
        "select",
        "option",
        "button",
        "input",
        "iframe",
        "audio",
        "video",
    }
)
"""Tags never exported: code, decorations and interactive shell widgets."""

SHELL_CLASSES = frozenset({"toolbar", "shell-header", "status-bar", "nav-arrow", "stage", "slide-frame"})
"""Classes of the navigation shell, outside of any slide content."""

BLOCK_TAGS = frozenset(
    {
        "div",
        "p",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "ul",
        "ol",
        "li",
        "table",
        "section",
        "article",
        "header",
        "footer",
        "figure",
        "blockquote",
        "pre",
        "hr",
    }
)

GRID_CLASSES = frozenset({"cds-grid", "stat-grid", "mention-wrap"})

TILE_CLASSES = frozenset({"cds-tile", "stat-card", "mention-pill"})

INLINE_BLOCK_CLASSES = frozenset({"arch-node-label", "gantt-sub"})
"""Inline classes the templates display as blocks (they get their own paragraph)."""

INLINE_BLOCK_TAGS = frozenset({"small"})
"""Inline tags used as sub-labels in the templates (``display: block``)."""

FILL_KEYS = ("background", "background-color")
BAR_KEYS = ("border-left-color", "border-left", "border-top-color", "border-top")

_PX_FLEX_BASIS_RE = re.compile(r"0\s+0\s+(\d+(?:\.\d+)?)px")
"""Matches the fixed-width shorthand (e.g. ``flex:0 0 116px``) of a custom Gantt label cell."""

_ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

_SHAPES = {
    "box": MSO_SHAPE.ROUNDED_RECTANGLE,
    "rect": MSO_SHAPE.RECTANGLE,
    "circle": MSO_SHAPE.OVAL,
    "oval": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "cylinder": MSO_SHAPE.CAN,
    "cloud": MSO_SHAPE.CLOUD,
}

_NOTIF_VARIANTS = {
    "success": ("DEFBE6", "24A148", "0E6027"),
    "warning": ("FDF6DD", "F1C21B", "8E6A00"),
    "error": ("FFF1F1", "DA1E28", "A2191F"),
}
"""``class`` variant of ``.cds-notification`` mapped to (fill, bar, title)."""

_ANNOT_CLASS_COLORS = {
    "annot-blue": ("003A8D", "FFFFFF"),
    "annot-orange": ("FBAE40", "4A3000"),
    "annot-coral": ("EC6962", "FFFFFF"),
}

# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


@dataclass
class ExportReport:
    """Outcome of an export: slide count plus every dropped or guessed item."""

    slide_count: int = 0
    warnings: list[str] = field(default_factory=list)
    charter: str | None = None

    def warn(self, message: str) -> None:
        """Record a dropped or approximated item.

        Callers are responsible for surfacing the list; the CLI prints it and
        exits non zero when no slide could be exported.
        """
        self.warnings.append(message)
        logger.debug("export pptx: %s", message)


def to_pptx(input_html: str, output_pptx: str) -> ExportReport:
    """Convert an HTML presentation to a PPTX file.

    Args:
        input_html: Path to the source HTML file.
        output_pptx: Destination path of the generated ``.pptx``.

    Returns:
        An :class:`ExportReport` with the number of slides written and the list
        of warnings (missing images, unsupported components, empty documents).
    """
    with trace_span("export.pptx", {"file.path": str(Path(input_html).resolve())}) as span:
        report = _build_pptx(input_html, output_pptx)
        span.set_attribute("slide.count", report.slide_count)
        span.set_attribute("warning.count", len(report.warnings))
        logger.info(
            "PPTX export: %s -> %s (%d slide(s), %d warning(s))",
            input_html,
            output_pptx,
            report.slide_count,
            len(report.warnings),
        )
        return report


def _build_pptx(input_html: str, output_pptx: str) -> ExportReport:
    """Do the actual HTML to PPTX conversion (see :func:`to_pptx`)."""
    source = Path(input_html)
    soup = BeautifulSoup(source.read_text(encoding="utf-8"), "html.parser")
    report = ExportReport()

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]

    articles = find_slides(soup)
    if not articles:
        report.warn(
            'Aucun element data-type="slide" ni article.slide trouve: '
            "le corps du document est exporte comme une slide unique."
        )
        body = soup.body or soup
        articles = [body] if isinstance(body, Tag) else []

    resolver = StyleResolver.from_soup(soup)
    context = _Context(
        res=resolver,
        theme=resolver.theme,
        base_dir=source.resolve().parent,
        report=report,
    )
    for index, article in enumerate(articles):
        slide = prs.slides.add_slide(blank_layout)
        _SlideBuilder(context, slide, article, index).build()

    report.slide_count = len(prs.slides)
    report.charter = resolver.theme.key
    out = Path(output_pptx)
    out.parent.mkdir(parents=True, exist_ok=True)
    with trace_span("file.write", {"file.path": str(out), "export.format": "pptx"}):
        prs.save(str(out))
    return report


def find_slides(soup: BeautifulSoup) -> list[Tag]:
    """Return the slide elements of a document, outermost ones only.

    Detection is attribute based (``data-type="slide"``, any tag name), with a
    fallback on ``article.slide`` / ``section.slide`` for templates that do not
    carry the attribute.
    """
    tagged = [el for el in soup.find_all(attrs={"data-type": "slide"}) if isinstance(el, Tag)]
    if not tagged:
        tagged = [el for el in soup.find_all(["article", "section"]) if isinstance(el, Tag) and has_class(el, "slide")]
    return [el for el in tagged if not _has_slide_ancestor(el, tagged)]


def _has_slide_ancestor(element: Tag, slides: list[Tag]) -> bool:
    """Tell whether ``element`` is nested inside another detected slide."""
    known = {id(s) for s in slides}
    parent = element.parent
    while parent is not None:
        if id(parent) in known:
            return True
        parent = parent.parent
    return False


# ---------------------------------------------------------------------------
# Internal model
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class _Context:
    """Shared state of one export run."""

    res: StyleResolver
    theme: Theme
    base_dir: Path
    report: ExportReport


@dataclass
class _Block:
    """A laid out piece of slide content, with its measured natural height."""

    kind: str
    element: Tag
    height: float
    min_height: float = 0.0
    max_height: float = 0.0
    flexible: bool = False
    style: TextStyle | None = None


@dataclass(frozen=True)
class _Region:
    """A rectangle of the slide that hosts a vertical flow of blocks."""

    box: Box
    elements: list[Tag]
    anchor: str = "top"


@dataclass(frozen=True)
class _Para:
    """One paragraph: a style, a nesting level and a list of styled runs."""

    style: TextStyle
    runs: list[tuple[str, TextStyle]]
    level: int = 0


# ---------------------------------------------------------------------------
# Slide builder
# ---------------------------------------------------------------------------


class _SlideBuilder:
    """Render one HTML slide element onto one python-pptx slide."""

    def __init__(self, context: _Context, slide: Any, article: Tag, index: int) -> None:
        self.ctx = context
        self.res = context.res
        self.theme = context.theme
        self.slide = slide
        self.article = article
        self.index = index
        self.handled: set[int] = set()

    # -- entry point -------------------------------------------------------

    def build(self) -> None:
        """Draw the chrome then flow the content of the slide."""
        kind = self._slide_kind()
        if self.theme.key == "ei":
            regions = self._chrome_ei(kind)
        elif self.theme.key == "carbon":
            regions = self._chrome_carbon()
        else:
            regions = self._chrome_generic()
        for region in regions:
            self._flow(region)

    def _slide_kind(self) -> str:
        """Return the ``data-slide-type`` of the slide, inferred when absent."""
        kind = str(self.article.get("data-slide-type") or "").strip()
        if kind:
            return kind
        if self.article.find(class_="slide-section-body") is not None:
            return "section"
        if self.article.find(class_="slide-cover-img") is not None:
            return "title"
        if self.article.find(class_="slide-inner") is not None:
            return "content"
        return ""

    # -- chrome ------------------------------------------------------------

    def _chrome_ei(self, kind: str) -> list[_Region]:
        """Draw the Euro-Information charter and return the content regions."""
        canvas = Box.slide()
        if kind == "section":
            self._rect(canvas, fill=self.theme.primary)
            band = self.article.find(class_="slide-section-band")
            if isinstance(band, Tag):
                self._consume(band)
            self._rect(canvas.pct(0, 72, 100, 6), fill=self.theme.primary_alt)
            body = self.article.find(class_="slide-section-body")
            host = body if isinstance(body, Tag) else self.article
            return [_Region(canvas.inset_px(72, 40), self._children(host), anchor="middle")]

        if kind == "title":
            return self._chrome_ei_title(canvas)

        # content / agenda / diagram: blue frame, white rounded inner area.
        self._rect(canvas, fill=self.theme.primary)
        inner = canvas.inset_px(10)
        self._rect(inner, fill="FFFFFF", radius_px=16)
        foot = self.article.find(class_="slide-foot")
        if isinstance(foot, Tag):
            self._render_ei_foot(foot)
        inner_host = self.article.find(class_="slide-inner")
        container = inner_host if isinstance(inner_host, Tag) else self.article
        if isinstance(inner_host, Tag):
            self._consume(inner_host, deep=False)
        content = inner.inset_px(40, 30, 40, 50)
        return [_Region(content, self._children(container))]

    def _chrome_ei_title(self, canvas: Box) -> list[_Region]:
        """Render the EI cover slide: full width image, centered title, logos."""
        cover = self.article.find(class_="slide-cover-img")
        if isinstance(cover, Tag):
            self._consume(cover)
            self._add_picture(canvas.pct(0, 0, 100, 66), str(cover.get("src") or ""), mode="cover")
        logos = self.article.find(class_="slide-cover-logos")
        if isinstance(logos, Tag):
            self._consume(logos)
            self._render_cover_logos(logos)
        body = self.article.find(class_="slide-cover-body")
        if isinstance(body, Tag):
            self._consume(body, deep=False)
            region = canvas.pct(0, 66, 100, 26).inset_px(48, 6)
            return [_Region(region, self._children(body), anchor="middle")]
        return [_Region(canvas.inset_px(48, 40), self._children(self.article), anchor="middle")]

    def _render_cover_logos(self, container: Tag) -> None:
        """Place the cover logo row: grouped logos left, single logo right."""
        baseline = (540.0 - 18.0) * PX_IN
        heights = {"logo-cm": 26.0, "logo-cic": 24.0, "logo-ei": 34.0}
        left_group = container.find(class_="logos-left")
        left_images = list(left_group.find_all("img")) if isinstance(left_group, Tag) else []
        right_images = [img for img in container.find_all("img") if img not in left_images and isinstance(img, Tag)]

        cursor = 48.0 * PX_IN
        for img in left_images:
            height = self._logo_height(img, heights) * PX_IN
            box = Box(cursor, baseline - height, 3.0, height)
            placed = self._add_picture(box, str(img.get("src") or ""), mode="fit-height")
            cursor += (inches(placed.width) if placed else 0.6) + 20.0 * PX_IN
        right_edge = SLIDE_W_IN - 48.0 * PX_IN
        for img in reversed(right_images):
            height = self._logo_height(img, heights) * PX_IN
            box = Box(right_edge - 1.2, baseline - height, 1.2, height)
            placed = self._add_picture(box, str(img.get("src") or ""), mode="fit-height")
            if placed:
                width = inches(placed.width)
                placed.left = Inches(right_edge - width)
                right_edge -= width + 20.0 * PX_IN

    @staticmethod
    def _logo_height(img: Tag, heights: dict[str, float]) -> float:
        """Return the CSS height of a cover logo from its class, default 24px."""
        for name in classes(img):
            if name in heights:
                return heights[name]
        return 24.0

    def _render_ei_foot(self, foot: Tag) -> None:
        """Render the EI footer: logo ring, page number and meeting title."""
        self._consume(foot)
        band = Box(0.0, SLIDE_H_IN - 40.0 * PX_IN, SLIDE_W_IN, 40.0 * PX_IN)
        ring = Box(2.0 * PX_IN, SLIDE_H_IN - 52.0 * PX_IN, 50.0 * PX_IN, 50.0 * PX_IN)
        oval = self._rect(ring, fill="FFFFFF", shape=MSO_SHAPE.OVAL)
        oval.line.color.rgb = RGBColor.from_string(self.theme.primary)
        oval.line.width = Pt(10)
        logo = foot.find("img")
        if isinstance(logo, Tag):
            inner = Box(ring.left, ring.top + (ring.height - 16.0 * PX_IN) / 2, ring.width, 16.0 * PX_IN)
            self._add_picture(inner, str(logo.get("src") or ""), mode="fit")

        page = foot.find(class_="slide-foot-page")
        if isinstance(page, Tag):
            box = Box(64.0 * PX_IN, band.top + 0.10, 0.8, 0.2)
            self._render_text(box, page, anchor="middle")
        title = foot.find(class_="slide-foot-title")
        if isinstance(title, Tag):
            box = Box(1.6, band.top + 0.10, SLIDE_W_IN - 3.2, 0.2)
            self._render_text(box, title, anchor="middle")

    def _chrome_carbon(self) -> list[_Region]:
        """Draw the IBM Carbon charter and return the content regions."""
        canvas = Box.slide()
        background = self._element_color(self.article) or "FFFFFF"
        self._rect(canvas, fill=background)
        regions: list[_Region] = []
        pad_x, pad_y = 40.0 * PX_IN, 30.0 * PX_IN
        content_w = SLIDE_W_IN - 2 * pad_x
        cursor = pad_y
        section = self._slide_kind() == "section"

        footer = self.article.find(class_="slide-footer")
        footer_top = SLIDE_H_IN
        if isinstance(footer, Tag):
            footer_top = SLIDE_H_IN - 40.0 * PX_IN
            self._rect(
                Box(0.0, footer_top, SLIDE_W_IN, 40.0 * PX_IN),
                fill=self._element_color(footer) or (self.theme.surface if is_light(background) else background),
                line=self.theme.border,
            )
            self._render_footer_texts(footer, footer_top)

        header = self.article.find(class_="slide-header")
        if isinstance(header, Tag):
            self._consume(header, deep=False)
            blocks = self._collect_all(self._children(header), content_w)
            height = sum(b.height for b in blocks)
            props = self.res.props(header)
            rule_h = parse_px(props.get("border-bottom", "")) * PX_IN or 4.0 * PX_IN
            rule_color = color_of(props, ("border-bottom-color", "border-bottom"), self.res)
            if section:
                # A section separator centers its title and rules off the footer.
                rule_y = footer_top - rule_h
                regions.append(
                    _Region(
                        Box(pad_x, pad_y, content_w, max(rule_y - 2 * pad_y, height)),
                        self._children(header),
                        anchor="middle",
                    )
                )
            else:
                rule_y = cursor + height + 14.0 * PX_IN
                regions.append(_Region(Box(pad_x, cursor, content_w, height), self._children(header)))
            self._rect(Box(0.0, rule_y, SLIDE_W_IN, rule_h), fill=rule_color or self.theme.primary)
            cursor = rule_y + rule_h + 22.0 * PX_IN

        body = self.article.find(class_="slide-body")
        host = body if isinstance(body, Tag) else self.article
        if isinstance(body, Tag):
            self._consume(body, deep=False)
        body_h = max(footer_top - 20.0 * PX_IN - cursor, 1.0)
        regions.append(_Region(Box(pad_x, cursor, content_w, body_h), self._children(host)))
        return regions

    def _render_footer_texts(self, footer: Tag, top: float) -> None:
        """Render the left and right texts of a Carbon style footer."""
        self._consume(footer)
        pad = 40.0 * PX_IN
        left = footer.find(class_="slide-footer-left")
        right = footer.find(class_="slide-footer-right")
        half = (SLIDE_W_IN - 2 * pad) / 2
        if isinstance(left, Tag):
            self._render_text(Box(pad, top + 0.09, half, 0.22), left, anchor="middle")
        if isinstance(right, Tag):
            self._render_text(Box(pad + half, top + 0.09, half, 0.22), right, anchor="middle")

    def _chrome_generic(self) -> list[_Region]:
        """Fallback layout for documents without a known charter."""
        canvas = Box.slide()
        self._rect(canvas, fill=self._element_color(self.article) or "FFFFFF")
        return [_Region(canvas.inset_px(48, 40), self._children(self.article))]

    # -- block flow --------------------------------------------------------

    def _flow(self, region: _Region) -> None:
        """Stack the blocks of a region, absorbing or shrinking extra space."""
        blocks = self._collect_all(region.elements, region.box.width)
        if not blocks:
            return
        available = region.box.height
        total = sum(b.height for b in blocks)
        flexible = [b for b in blocks if b.flexible]

        if total < available and flexible:
            free = available - total
            for index, block in enumerate(flexible):
                share = free / (len(flexible) - index)
                room = block.max_height - block.height if block.max_height else share
                grow = max(min(share, room), 0.0)
                block.height += grow
                free -= grow
        elif total > available:
            overflow = total - available
            for block in flexible:
                room = max(block.height - block.min_height, 0.0)
                take = min(room, overflow)
                block.height -= take
                overflow -= take
                if overflow <= 0.001:
                    break
            if overflow > 0.001:
                scale = available / (available + overflow)
                for block in blocks:
                    block.height *= scale

        cursor = region.box.top
        if region.anchor == "middle":
            used = sum(b.height for b in blocks)
            cursor += max((region.box.height - used) / 2.0, 0.0)
        for block in blocks:
            box = Box(region.box.left, cursor, region.box.width, block.height)
            self._render_block(block, box)
            cursor += block.height

    def _collect_all(self, elements: list[Tag], width: float) -> list[_Block]:
        """Collect the blocks of a list of sibling elements."""
        for element in elements:
            if isinstance(element, Tag) and self._is_inline_gantt(element):
                self._consume_inline_gantt_siblings(element)
        blocks: list[_Block] = []
        for element in elements:
            blocks.extend(self._collect(element, width))
        return blocks

    def _consume_inline_gantt_siblings(self, element: Tag) -> None:
        """Mark a custom inline Gantt's header/legend siblings as already rendered.

        Both live outside ``element`` (see :meth:`_inline_gantt_sibling`), so
        without this a sibling-order walk renders the header row as one huge
        text block (it comes first in the DOM) before ``_gantt_block`` ever
        gets a chance to claim it.
        """
        legend = self._inline_gantt_legend(element)
        if legend is not None:
            self._consume(legend)
        head = self._inline_gantt_sibling(element, lambda row: parse_px(style_props(row).get("margin-left", "")) > 0.0)
        if head is not None:
            self._consume(head)

    def _collect(self, element: Tag, width: float) -> list[_Block]:
        """Turn one element into zero, one or several layout blocks."""
        if not isinstance(element, Tag) or element.name in SKIP_TAGS:
            return []
        if id(element) in self.handled or has_class(element, *SHELL_CLASSES):
            return []
        if str(element.get("aria-hidden") or "") == "true":
            return []

        dtype = str(element.get("data-type") or "")
        if dtype == "gantt" or self._is_inline_gantt(element):
            return [self._gantt_block(element, width)]
        if dtype == "arch-diagram":
            return [self._arch_block(element, width)]
        if dtype == "annotated-image":
            return [self._annotated_block(element, width)]
        if dtype == "table" or element.name == "table":
            return [self._table_block(element, width)]
        if has_class(element, "cds-notification"):
            return [self._panel_block(element, width)]
        if has_class(element, *GRID_CLASSES):
            return [self._grid_block(element, width)]
        if has_class(element, "cw-bar"):
            return [_Block("hbar", element, 32.0 * PX_IN + 0.08)]
        if has_class(element, "slide-title-rule", "accent-line"):
            return [_Block("rule", element, 33.0 * PX_IN)]
        if element.name == "img":
            return [self._image_block(element, width)]
        if element.name == "hr":
            return [_Block("rule", element, 12.0 * PX_IN)]
        if self._is_panel(element):
            return [self._panel_block(element, width)]
        if element.name in {"ul", "ol"} or not self._has_block_children(element):
            return self._text_block(element, width)
        return self._collect_all(self._children(element), width)

    def _is_panel(self, element: Tag) -> bool:
        """Tell whether an element is a decorated panel (fill or accent bar)."""
        if has_class(element, *TILE_CLASSES):
            return True
        props = self.res.props(element)
        return bool(color_of(props, FILL_KEYS, self.res) or color_of(props, BAR_KEYS, self.res))

    def _has_block_children(self, element: Tag) -> bool:
        """Tell whether an element contains children that become their own block."""
        return any(isinstance(child, Tag) and self._is_block_child(child) for child in element.children)

    @staticmethod
    def _is_block_child(child: Tag) -> bool:
        """Tell whether a child element is laid out as a block of its own."""
        if child.name in SKIP_TAGS:
            return False
        return (
            child.name in BLOCK_TAGS
            or child.name in INLINE_BLOCK_TAGS
            or child.name == "img"
            or bool(child.get("data-type"))
            or has_class(child, *INLINE_BLOCK_CLASSES, *GRID_CLASSES, *TILE_CLASSES)
        )

    @staticmethod
    def _children(element: Tag) -> list[Tag]:
        """Return the direct element children of a node."""
        return [child for child in element.children if isinstance(child, Tag)]

    def _consume(self, element: Tag, deep: bool = True) -> None:
        """Mark an element (and optionally its subtree) as already rendered."""
        self.handled.add(id(element))
        if deep:
            for node in element.find_all(True):
                self.handled.add(id(node))

    # -- block measurement -------------------------------------------------

    def _text_block(self, element: Tag, width: float) -> list[_Block]:
        """Measure a text block, or drop it when it carries no text."""
        style = self.res.style(element)
        paras = self._paragraphs(element, style)
        if not paras:
            return []
        height = sum(self._para_height(p, width) for p in paras)
        return [_Block("text", element, height, style=style)]

    def _para_height(self, para: _Para, width: float) -> float:
        """Height of one paragraph, indentation included."""
        text = "".join(text for text, _ in para.runs)
        indent = para.level * 0.18
        return block_height(text, para.style, max(width - indent, 0.5))

    def _text_height(self, element: Tag, width: float, style: TextStyle) -> float:
        """Height of the whole text content of an element."""
        return sum(self._para_height(p, width) for p in self._paragraphs(element, style))

    def _panel_block(self, element: Tag, width: float) -> _Block:
        """Measure a decorated panel: text plus its padding and accent bar."""
        style = self.res.style(element)
        height = self._text_height(element, width - 0.32, style) + 0.3
        return _Block("panel", element, max(height, 0.4))

    def _grid_block(self, element: Tag, width: float) -> _Block:
        """Measure a tile grid: rows of equal height cells."""
        cells = self._children(element)
        if not cells:
            return _Block("grid", element, 0.05)
        cols = self._grid_columns(element, len(cells))
        gap = 14.0 * PX_IN
        cell_w = max((width - gap * (cols - 1)) / cols, 0.4)
        cell_h = 0.0
        for cell in cells:
            style = self.res.style(cell)
            cell_h = max(cell_h, self._text_height(cell, cell_w - 0.28, style) + 0.24)
        rows = math.ceil(len(cells) / cols)
        natural = rows * cell_h + (rows - 1) * gap + 0.1
        return _Block(
            "grid",
            element,
            natural,
            min_height=natural * 0.85,
            max_height=natural * 1.7,
            flexible=True,
        )

    @staticmethod
    def _grid_columns(element: Tag, count: int) -> int:
        """Number of columns of a grid, from its ``cols-N`` class."""
        for name in classes(element):
            match = re.fullmatch(r"cols-(\d)", name)
            if match:
                return int(match.group(1))
        if has_class(element, "stat-grid"):
            return 3
        if has_class(element, "mention-wrap"):
            return min(3, max(count, 1))
        return min(count, 3) or 1

    def _table_block(self, element: Tag, width: float) -> _Block:
        """Measure a table from its own cell contents."""
        grid = TableGrid(element)
        if not grid.rows:
            return _Block("text", element, 0.05)
        widths = grid.column_widths(width)
        height = 0.0
        for row_index, row in enumerate(grid.rows):
            # PPTX rows never shrink below their text, so keep a safety margin.
            row_h = 0.3 if row_index in grid.header_rows else 0.34
            for cell in row:
                if cell is None or cell.origin != row_index:
                    continue
                style = self.res.style(cell.element)
                cell_w = sum(widths[cell.column : cell.column + cell.colspan]) - 0.18
                row_h = max(row_h, self._text_height(cell.element, max(cell_w, 0.3), style) + 0.16)
            grid.row_heights.append(row_h)
            height += row_h
        natural = height + 0.06
        return _Block(
            "table",
            element,
            natural,
            min_height=natural,
            max_height=natural * 1.6,
            flexible=True,
        )

    def _gantt_block(self, element: Tag, width: float) -> _Block:
        """Measure a Gantt chart: one line per task, plus scale and legend."""
        rows = self._gantt_rows(element)
        row_h_px = sum((row.track_h_px or 30.0) for row in rows) or 30.0 * max(len(rows), 1)
        height = row_h_px * PX_IN + 24.0 * PX_IN
        legend = element.find(class_="gantt-legend") or self._inline_gantt_legend(element)
        if isinstance(legend, Tag):
            height += self._gantt_legend_height(legend, width)
            self._consume(legend)
        inline_head = self._inline_gantt_sibling(
            element, lambda row: parse_px(style_props(row).get("margin-left", "")) > 0.0
        )
        if inline_head is not None:
            self._consume(inline_head)
        return _Block("gantt", element, height, min_height=height * 0.6)

    def _gantt_legend_height(self, legend: Tag, width: float) -> float:
        """Measure the legend band height, counting wrapped lines on a busy legend."""
        entries = [entry for entry in self._children(legend) if entry.get_text(strip=True)]
        if not entries:
            return 26.0 * PX_IN
        style = TextStyle(size=9.0, color="secondary", space_after=0)
        swatch_entries = [entry for entry in entries if not self._is_marker_legend_entry(entry)]
        marker_entries = [entry for entry in entries if self._is_marker_legend_entry(entry)]
        groups = [group for group in (swatch_entries, marker_entries) if group]
        line_count = sum(len(self._wrap_legend_entries(group, width, style)) for group in groups) or 1
        return line_count * 20.0 * PX_IN

    def _arch_block(self, element: Tag, width: float) -> _Block:
        """Measure a diagram: fixed when the CSS gives a height, else flexible."""
        props = style_props(element)
        explicit = parse_length(props.get("height"), 0.0)
        if explicit and explicit > 0.2:
            return _Block("arch", element, explicit, min_height=explicit * 0.7)
        return _Block("arch", element, 2.6, min_height=1.6, flexible=True)

    def _annotated_block(self, element: Tag, width: float) -> _Block:
        """Measure an annotated image from the aspect ratio of its picture."""
        img = element.find("img")
        aspect = 16.0 / 9.0
        if isinstance(img, Tag):
            measured = self._image_aspect(str(img.get("src") or ""))
            if measured:
                aspect = measured
        img_w = min(width, self._max_width(element, width))
        return _Block("annotated-image", element, img_w / aspect, min_height=1.0)

    def _image_block(self, element: Tag, width: float) -> _Block:
        """Measure a standalone image."""
        props = style_props(element)
        aspect = self._image_aspect(str(element.get("src") or "")) or 16.0 / 9.0
        height = parse_length(props.get("height"), SLIDE_H_IN)
        img_w = parse_length(props.get("width"), width) or width
        img_w = min(img_w, width)
        return _Block("image", element, height or img_w / aspect)

    @staticmethod
    def _max_width(element: Tag, width: float) -> float:
        """Resolve the ``max-width`` of an element, defaulting to ``width``."""
        value = parse_length(style_props(element).get("max-width"), width)
        return value if value else width

    # -- block rendering ---------------------------------------------------

    def _render_block(self, block: _Block, box: Box) -> None:
        """Dispatch a measured block to its renderer."""
        if block.kind == "text":
            self._render_text(box, block.element, style=block.style)
        elif block.kind == "rule":
            self._render_rule(box, block.element)
        elif block.kind == "panel":
            self._render_panel(box, block.element)
        elif block.kind == "hbar":
            self._render_hbar(box, block.element)
        elif block.kind == "grid":
            self._render_grid(box, block.element)
        elif block.kind == "table":
            self._render_table(box, block.element)
        elif block.kind == "gantt":
            self._render_gantt(box, block.element)
        elif block.kind == "arch":
            self._render_arch(box, block.element)
        elif block.kind == "annotated-image":
            self._render_annotated_image(box, block.element)
        elif block.kind == "image":
            self._render_image(box, block.element)
        else:  # pragma: no cover - defensive, every kind is handled above
            self.ctx.report.warn(f"Composant non supporte ignore: {block.kind}")

    def _render_text(
        self,
        box: Box,
        element: Tag,
        style: TextStyle | None = None,
        anchor: str = "top",
    ) -> None:
        """Render the text content of an element inside ``box``."""
        base = style or self.res.style(element)
        paras = self._paragraphs(element, base)
        if not paras:
            return
        _, frame = self._textbox(box, anchor=anchor)
        self._write(frame, paras)

    def _render_rule(self, box: Box, element: Tag) -> None:
        """Render the short accent rule found under EI slide titles."""
        props = self.res.props(element)
        width = parse_length(props.get("width"), box.width) or 64.0 * PX_IN
        height = parse_length(props.get("height"), 0.0) or 3.0 * PX_IN
        color = color_of(props, FILL_KEYS, self.res)
        rule = Box(box.left, box.top + 12.0 * PX_IN, min(width, box.width), height)
        self._rect(rule, fill=color or self.theme.accent)

    def _render_panel(self, box: Box, element: Tag) -> None:
        """Render a callout or card: tinted panel, accent bar, then its text."""
        fill, bar, vertical = self._panel_colors(element)
        panel = Box(box.left, box.top, box.width, max(box.height - 0.06, 0.2))
        self._rect(panel, fill=fill, line=None if fill else self.theme.border)
        if bar and vertical:
            self._rect(Box(panel.left, panel.top, 4.0 * PX_IN, panel.height), fill=bar)
        elif bar:
            self._rect(Box(panel.left, panel.top, panel.width, 3.0 * PX_IN), fill=bar)
        inner = Box(panel.left + 0.17, panel.top + 0.08, panel.width - 0.3, panel.height - 0.14)
        if self._has_block_children(element) and any(
            self._is_panel(child) or child.name == "table" for child in self._children(element)
        ):
            self._flow(_Region(inner, self._children(element)))
        else:
            self._render_text(inner, element)

    def _panel_colors(self, element: Tag) -> tuple[str | None, str | None, bool]:
        """Resolve the (fill, accent bar, bar is vertical) look of a panel."""
        props = self.res.props(element)
        fill = color_of(props, FILL_KEYS, self.res)
        left_bar = color_of(props, ("border-left-color", "border-left"), self.res)
        top_bar = color_of(props, ("border-top-color", "border-top"), self.res)
        for variant, colors in _NOTIF_VARIANTS.items():
            if has_class(element, "cds-notification") and has_class(element, variant):
                return colors[0], colors[1], True
        if has_class(element, "cds-notification"):
            default_fill = "FFF6E6" if self.theme.key == "ei" else "EDF5FF"
            accent = self.theme.accent if self.theme.key == "ei" else self.theme.primary
            return fill or default_fill, left_bar or accent, True
        if left_bar:
            return fill or self.theme.surface, left_bar, True
        if top_bar:
            return fill or self.theme.surface, top_bar, False
        return fill or self.theme.surface, None, True

    def _render_hbar(self, box: Box, element: Tag) -> None:
        """Render a proportional horizontal bar such as a context window chart."""
        segments = [seg for seg in self._children(element) if seg.name not in SKIP_TAGS]
        if not segments:
            return
        weights: list[float] = []
        for segment in segments:
            declared = parse_pct(style_props(segment).get("width", ""))
            weights.append(declared if declared > 0 else 0.0)
        known = sum(weights)
        unknown = [i for i, w in enumerate(weights) if w == 0.0]
        for index in unknown:
            weights[index] = max((100.0 - known) / len(unknown), 2.0)
        total = sum(weights) or 100.0
        cursor = box.left
        bar = Box(box.left, box.top, box.width, max(box.height - 0.08, 0.16))
        for segment, weight in zip(segments, weights, strict=True):
            width = bar.width * weight / total
            fill = self._element_color(segment) or self.theme.primary
            cell = Box(cursor, bar.top, max(width, 0.04), bar.height)
            self._rect(cell, fill=fill)
            text = segment.get_text(" ", strip=True)
            if text and width > 0.3:
                style = replace(self.res.style(segment), size=9.0, bold=True, align="center", space_after=0)
                _, frame = self._textbox(cell, anchor="middle")
                self._write(frame, [_Para(style, [(text, style)])])
            cursor += width

    def _render_grid(self, box: Box, element: Tag) -> None:
        """Render a grid of tiles, stat cards or pills."""
        cells = self._children(element)
        if not cells:
            return
        cols = self._grid_columns(element, len(cells))
        gap = 14.0 * PX_IN
        rows = math.ceil(len(cells) / cols)
        cell_w = max((box.width - gap * (cols - 1)) / cols, 0.4)
        cell_h = max((box.height - gap * (rows - 1) - 0.1) / rows, 0.3)
        for position, cell in enumerate(cells):
            row, col = divmod(position, cols)
            cell_box = Box(
                box.left + col * (cell_w + gap),
                box.top + row * (cell_h + gap),
                cell_w,
                cell_h,
            )
            self._render_cell(cell_box, cell)

    def _render_cell(self, box: Box, cell: Tag) -> None:
        """Render one grid cell: decorated tile, nested blocks or plain text."""
        if has_class(cell, *TILE_CLASSES) or self._is_panel(cell):
            self._render_panel(box, cell)
            return
        if self._has_block_children(cell):
            self._flow(_Region(box, self._children(cell)))
            return
        self._render_text(box, cell)

    def _render_table(self, box: Box, element: Tag) -> None:
        """Render an HTML table as a native PPTX table with merges and styling."""
        grid = TableGrid(element)
        if not grid.rows:
            return
        widths = grid.column_widths(box.width)
        shape = self.slide.shapes.add_table(
            len(grid.rows), grid.column_count, *Box(box.left, box.top, box.width, box.height).emu()
        )
        table = shape.table
        table.first_row = bool(grid.header_rows)
        table.horz_banding = False
        for index, width in enumerate(widths):
            table.columns[index].width = Inches(width)
        heights = grid.row_heights or [box.height / len(grid.rows)] * len(grid.rows)
        scale = box.height / max(sum(heights), 0.01)
        for index, height in enumerate(heights):
            table.rows[index].height = Inches(height * scale)

        striped = str(element.get("data-style") or "") == "striped"
        for r_index, row in enumerate(grid.rows):
            is_header = r_index in grid.header_rows
            for c_index, cell in enumerate(row):
                if cell is None or cell.origin != r_index or cell.column != c_index:
                    continue
                target = table.cell(r_index, c_index)
                if cell.colspan > 1 or cell.rowspan > 1:
                    end_row = min(r_index + cell.rowspan - 1, len(grid.rows) - 1)
                    end_col = min(c_index + cell.colspan - 1, grid.column_count - 1)
                    target.merge(table.cell(end_row, end_col))
                self._fill_table_cell(target, cell.element, is_header, striped and r_index % 2 == 0)
                if not is_header and r_index < len(grid.rows) - 1:
                    set_cell_border(target, "B", self.theme.border)

    def _fill_table_cell(self, target: Any, element: Tag, header: bool, striped: bool) -> None:
        """Write and style one table cell."""
        target.margin_left = Inches(0.07)
        target.margin_right = Inches(0.07)
        target.margin_top = Inches(0.03)
        target.margin_bottom = Inches(0.03)
        target.vertical_anchor = MSO_ANCHOR.MIDDLE
        if header:
            target.fill.solid()
            target.fill.fore_color.rgb = RGBColor.from_string(self.theme.table_header_bg)
        elif striped:
            target.fill.solid()
            target.fill.fore_color.rgb = RGBColor.from_string(self.theme.surface)
        else:
            target.fill.background()
        default = self.res.style(element)
        if header:
            default = replace(default, color=self.theme.table_header_fg, bold=True)
        paras = self._paragraphs(element, default)
        frame = target.text_frame
        frame.word_wrap = True
        self._write(frame, paras or [_Para(default, [("", default)])])

    def _render_gantt(self, box: Box, element: Tag) -> None:
        """Render a Gantt chart as real bars on a quarter ruled track."""
        rows = self._gantt_rows(element)
        legend = element.find(class_="gantt-legend") or self._inline_gantt_legend(element)
        legend_h = self._gantt_legend_height(legend, box.width) if isinstance(legend, Tag) else 0.0
        head_h = 24.0 * PX_IN
        label_w = min(GANTT_LABEL_PX * PX_IN, box.width * 0.3)
        dates_w = (
            min(GANTT_DATES_PX * PX_IN, box.width * 0.16) if element.find(class_="gantt-dates") is not None else 0.0
        )
        track_x = box.left + label_w
        track_w = max(box.width - label_w - dates_w, 0.5)
        rows_h = max(box.height - head_h - legend_h, 0.3)
        natural_h = [row.track_h_px or 30.0 for row in rows] or [30.0]
        scale = rows_h / (sum(natural_h) * PX_IN)
        row_heights = [h * PX_IN * scale for h in natural_h]

        self._render_gantt_head(element, Box(track_x, box.top, track_w, head_h))
        track_top = box.top + head_h
        self._rect(
            Box(track_x, track_top, track_w, rows_h),
            fill=self.theme.surface,
            line=self.theme.border,
        )
        for index in (25.0, 50.0, 75.0):
            x = track_x + track_w * index / 100.0
            self._rect(Box(x, track_top, 1.0 * PX_IN, rows_h), fill=self.theme.border)

        period = gantt_period(element, rows)
        top = track_top
        for row, row_h in zip(rows, row_heights, strict=True):
            if row.label is not None:
                self._render_text(Box(box.left, top, label_w - 0.08, row_h), row.label, anchor="middle")
            for marker in row.markers:
                marker_left_pct = parse_pct(style_props(marker).get("left", ""))
                marker_color = self._element_color(marker) or self.theme.border
                self._rect(
                    Box(track_x + track_w * marker_left_pct / 100.0, top, 1.2 * PX_IN, row_h),
                    fill=marker_color,
                )
            for task in row.tasks:
                left_pct, width_pct = gantt_geometry(task, period)
                top_pct, height_pct = gantt_task_band(task, row.track_h_px)
                bar = Box(
                    track_x + track_w * left_pct / 100.0,
                    top + row_h * top_pct / 100.0,
                    max(track_w * width_pct / 100.0, 0.12),
                    max(row_h * height_pct / 100.0, 0.06),
                )
                color = self._element_color(task) or self.theme.primary
                hatched = "repeating-linear-gradient" in style_props(task).get("background-image", "")
                self._rect(bar, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius_px=3, hatch=hatched)
                text = task.get_text(" ", strip=True) or str(task.get("data-label") or "")
                if text:
                    label, style = self._fit_bar_label(text, bar)
                    _, frame = self._textbox(bar.inset(0.06, 0.0), anchor="middle", wrap=False)
                    self._write(frame, [_Para(style, [(label, style)])])
            if row.dates is not None:
                self._render_text(
                    Box(track_x + track_w + 0.06, top, max(dates_w - 0.06, 0.3), row_h),
                    row.dates,
                    anchor="middle",
                )
            top += row_h
        if isinstance(legend, Tag):
            self._render_gantt_legend(Box(box.left, box.bottom - legend_h, box.width, legend_h), legend)

    @staticmethod
    def _fit_bar_label(text: str, bar: Box) -> tuple[str, TextStyle]:
        """Pick a font size (then truncate) so a bar label stays on one line.

        A short sub-lane bar in a dense roadmap is often narrower than its
        full label at any readable size (the source CSS relies on the
        browser's own text overflow clipping, which PPTX has no equivalent
        of): shrink down to a 5.5pt floor first, matching the template's own
        ``font-size:5.5px`` convention for these bars, then truncate with an
        ellipsis rather than let python-pptx's ``word_wrap`` push the
        overflow outside the shape, which renders as illegible ghost text
        over the neighbouring rows.
        """
        usable_w_pt = max((bar.width - 0.12) * 72.0, 4.0)
        for size in (9.5, 8.0, 6.5, 5.5):
            if len(text) * size * 0.5 <= usable_w_pt:
                return text, TextStyle(size=size, bold=True, color="FFFFFF", space_after=0)
        size = 5.5
        max_chars = max(int(usable_w_pt / (size * 0.5)), 1)
        if max_chars >= len(text):
            return text, TextStyle(size=size, bold=True, color="FFFFFF", space_after=0)
        truncated = text[: max(max_chars - 1, 1)].rstrip() + "…"
        return truncated, TextStyle(size=size, bold=True, color="FFFFFF", space_after=0)

    def _render_gantt_head(self, element: Tag, box: Box) -> None:
        """Render the period scale of a Gantt chart above the track."""
        head = element.find(class_="gantt-head") or element.find(class_="gantt-scale")
        if isinstance(head, Tag):
            track = head.find(class_="gantt-track")
            cells = (
                self._children(track)
                if isinstance(track, Tag)
                else [cell for cell in self._children(head) if not has_class(cell, "gantt-label-col")]
            )
            self._render_gantt_head_cells(box, [cell for cell in cells if cell.get_text(strip=True)])
            return
        inline_head = self._inline_gantt_sibling(
            element, lambda row: parse_px(style_props(row).get("margin-left", "")) > 0.0
        )
        if isinstance(inline_head, Tag):
            cells = [cell for cell in self._children(inline_head) if cell.get_text(strip=True)]
            self._render_gantt_head_cells(box, cells)

    def _render_gantt_head_cells(self, box: Box, cells: list[Tag]) -> None:
        """Lay out the period cells of a Gantt head at equal width across ``box``."""
        if not cells:
            return
        cell_w = box.width / len(cells)
        for index, cell in enumerate(cells):
            style = replace(
                self.res.style(cell),
                align="center",
                size=10.5,
                bold=True,
                color="secondary",
            )
            self._render_text(
                Box(box.left + index * cell_w, box.top, cell_w, box.height),
                cell,
                style=style,
                anchor="middle",
            )

    @staticmethod
    def _inline_gantt_sibling(element: Tag, predicate: Callable[[Tag], bool]) -> Tag | None:
        """Find the sibling of a custom inline Gantt's row stack matching ``predicate``.

        The header row (month cells) and the legend row sit as siblings of the
        row-stack container inside their common parent (e.g. ``.slide-body``),
        not as children of it: the row stack is only ``flex:1;overflow-y:auto``
        wrapping the lanes, so both must be looked up on the parent.
        """
        parent = element.parent
        if not isinstance(parent, Tag):
            return None
        for sibling in parent.find_all(True, recursive=False):
            if isinstance(sibling, Tag) and sibling is not element and predicate(sibling):
                return sibling
        return None

    def _render_gantt_legend(self, box: Box, legend: Tag) -> None:
        """Render the color legend of a Gantt chart: swatches, then marker keys.

        Two entry kinds coexist in both the documented ``gantt-legend`` shape
        and the custom inline one: a color-swatch entry (lane category) and a
        marker-key entry (a thin vertical bar, matching a milestone line drawn
        in the track). The HTML visually separates the marker-key group with a
        leading border; entries are laid out in two left-aligned rows so that
        separation survives the export instead of forcing every entry into one
        equal-width row (which would also make a long label collide with the
        next swatch on a busy legend such as the client's 15-entry one).
        """
        entries = [entry for entry in self._children(legend) if entry.get_text(strip=True)]
        if not entries:
            return
        swatch_entries = [entry for entry in entries if not self._is_marker_legend_entry(entry)]
        marker_entries = [entry for entry in entries if self._is_marker_legend_entry(entry)]
        groups = [group for group in (swatch_entries, marker_entries) if group]
        style = TextStyle(size=9.0, color="secondary", space_after=0)
        wrapped = [self._wrap_legend_entries(group, box.width, style) for group in groups]
        line_count = sum(len(lines) for lines in wrapped) or 1
        line_h = box.height / line_count
        top = box.top
        for lines in wrapped:
            for line in lines:
                self._render_gantt_legend_row(Box(box.left, top, box.width, line_h), line, style)
                top += line_h

    def _wrap_legend_entries(self, entries: list[Tag], width: float, style: TextStyle) -> list[list[Tag]]:
        """Split legend entries into lines that fit ``width``, greedy left to right."""
        lines: list[list[Tag]] = [[]]
        cursor = 0.0
        for entry in entries:
            entry_w = self._legend_entry_width(entry, style)
            if cursor + entry_w > width and lines[-1]:
                lines.append([])
                cursor = 0.0
            lines[-1].append(entry)
            cursor += entry_w
        return [line for line in lines if line]

    def _legend_entry_width(self, entry: Tag, style: TextStyle) -> float:
        """Estimate the rendered width of one legend entry (swatch plus label)."""
        swatch_w = 0.0
        if self._legend_swatch(entry) is not None:
            swatch_w = (3.0 if self._is_marker_legend_entry(entry) else 11.0) * PX_IN + 0.06
        text = entry.get_text(" ", strip=True)
        text_w = len(text) * style.size * PX_IN * 1.05
        return swatch_w + text_w + 0.2

    def _render_gantt_legend_row(self, box: Box, entries: list[Tag], style: TextStyle) -> None:
        """Render one line of legend entries at their natural (unequal) width."""
        cursor = box.left
        gap = 0.16
        for entry in entries:
            swatch = self._legend_swatch(entry)
            text = entry.get_text(" ", strip=True)
            if isinstance(swatch, Tag):
                color = self._element_color(swatch) or self.theme.primary
                marker = self._is_marker_legend_entry(entry)
                swatch_w = 3.0 * PX_IN if marker else 11.0 * PX_IN
                hatched = "repeating-linear-gradient" in style_props(swatch).get("background-image", "")
                self._rect(
                    Box(cursor, box.top + (box.height - 11.0 * PX_IN) / 2.0, swatch_w, 11.0 * PX_IN),
                    fill=color,
                    hatch=hatched,
                )
                swatch.extract()
                cursor += swatch_w + 0.06
            text_w = max(len(text) * style.size * PX_IN * 1.05, 0.3)
            self._render_text(Box(cursor, box.top, text_w, box.height), entry, style=style, anchor="middle")
            cursor += text_w + gap

    @staticmethod
    def _legend_swatch(entry: Tag) -> Tag | None:
        """Return the swatch element of a legend entry: ``<i>`` or a swatch ``<span>``."""
        swatch = entry.find("i")
        if isinstance(swatch, Tag):
            return swatch
        for span in entry.find_all("span", recursive=False):
            if isinstance(span, Tag) and not span.get_text(strip=True):
                return span
        return None

    def _is_marker_legend_entry(self, entry: Tag) -> bool:
        """Tell a milestone marker-key entry (thin vertical bar) from a color swatch.

        The HTML draws a marker key as a 1-2px wide, taller bar instead of the
        wider, shorter rounded rectangle used for lane categories; a leading
        ``border-left`` on the entry itself is the other tell used in the
        client deck, kept as a secondary signal since a future author could
        drop it.
        """
        swatch = self._legend_swatch(entry)
        if isinstance(swatch, Tag):
            props = style_props(swatch)
            width_px = parse_px(props.get("width", ""))
            height_px = parse_px(props.get("height", ""))
            if width_px and height_px:
                return width_px <= 3.0 and height_px > width_px
        return "border-left" in style_props(entry)

    def _inline_gantt_legend(self, element: Tag) -> Tag | None:
        """Find the legend row of a custom inline Gantt, if any.

        The legend is not marked with ``gantt-legend``: it is the sibling flex
        row after the row stack (see :meth:`_inline_gantt_sibling`), wrapping
        (``flex-wrap:wrap``) entries that are themselves flex rows with a
        swatch ``<span>`` plus a label.
        """

        def _is_legend(candidate: Tag) -> bool:
            if style_props(candidate).get("flex-wrap") != "wrap":
                return False
            entries = self._children(candidate)
            return bool(entries) and all(self._legend_swatch(entry) is not None for entry in entries)

        return self._inline_gantt_sibling(element, _is_legend)

    def _is_inline_gantt(self, element: Tag) -> bool:
        """Structurally detect a custom inline-styled Gantt (no CSS classes).

        Recognizes a container whose direct children are at least two row-like
        divs (see :meth:`_is_inline_gantt_row`): a fixed-width label cell
        followed by a ``position:relative`` track sibling holding absolutely
        positioned bars and/or milestone marker lines. Structural, not
        attribute-based, so it generalizes to any hand-authored deck sharing
        the shape without per-file markup changes.
        """
        if str(element.get("data-type") or "") or element.find(class_="gantt-row") is not None:
            return False
        rows = [child for child in self._children(element) if self._is_inline_gantt_row(child)]
        return len(rows) >= 2

    @staticmethod
    def _is_inline_gantt_row(row: Tag) -> bool:
        """Tell whether a div is one lane of a custom inline Gantt.

        Shape: a flex row (``display:flex``) whose first element child has a
        fixed ``flex-basis``/width (the label cell) and whose second element
        child is ``position:relative`` with ``flex:1`` (the track) and
        contains at least one absolutely positioned descendant (a bar or a
        marker line).
        """
        if row.name != "div":
            return False
        row_props = style_props(row)
        if row_props.get("display") != "flex":
            return False
        cells = [c for c in row.children if isinstance(c, Tag)]
        if len(cells) < 2:
            return False
        label_cell, track = cells[0], cells[1]
        label_props = style_props(label_cell)
        has_fixed_label = bool(_PX_FLEX_BASIS_RE.search(label_props.get("flex", ""))) or bool(
            parse_px(label_props.get("width", ""))
        )
        track_props = style_props(track)
        is_track = track_props.get("position") == "relative" and "1" in track_props.get("flex", "")
        if not (has_fixed_label and is_track):
            return False
        return any(style_props(c).get("position") == "absolute" for c in track.children if isinstance(c, Tag))

    def _inline_gantt_rows(self, element: Tag) -> list[GanttRow]:
        """Extract rows, tasks and milestone markers of a custom inline Gantt."""
        rows: list[GanttRow] = []
        for row in self._children(element):
            if not self._is_inline_gantt_row(row):
                continue
            cells = [c for c in row.children if isinstance(c, Tag)]
            label_cell, track = cells[0], cells[1]
            track_h_px = parse_px(style_props(row).get("min-height", "")) or parse_px(
                style_props(track).get("min-height", "")
            )
            tasks: list[Tag] = []
            markers: list[Tag] = []
            for child in track.children:
                if not isinstance(child, Tag) or style_props(child).get("position") != "absolute":
                    continue
                if child.get_text(strip=True):
                    tasks.append(child)
                else:
                    markers.append(child)
            rows.append(GanttRow(label=label_cell, dates=None, tasks=tasks, markers=markers, track_h_px=track_h_px))
        return rows

    def _gantt_rows(self, element: Tag) -> list[GanttRow]:
        """Extract the rows of a Gantt chart, with a flat task fallback."""
        rows: list[GanttRow] = []
        for row in element.find_all(class_="gantt-row"):
            tasks = [task for task in row.find_all(attrs={"data-type": "gantt-task"}) if isinstance(task, Tag)]
            label = row.find(class_="gantt-label")
            dates = row.find(class_="gantt-dates")
            rows.append(
                GanttRow(
                    label=label if isinstance(label, Tag) else None,
                    dates=dates if isinstance(dates, Tag) else None,
                    tasks=tasks,
                )
            )
        if rows:
            return rows
        rows = self._inline_gantt_rows(element)
        if rows:
            return rows
        for task in element.find_all(attrs={"data-type": "gantt-task"}):
            if isinstance(task, Tag):
                rows.append(GanttRow(label=None, dates=None, tasks=[task]))
        if not rows:
            self.ctx.report.warn('Gantt sans tache data-type="gantt-task": bloc ignore.')
        return rows

    def _render_arch(self, box: Box, element: Tag) -> None:
        """Render an architecture diagram: nodes as shapes, edges as lines."""
        background = self._element_color(element)
        if background:
            self._rect(box, fill=background, line=self.theme.border)
        nodes = 0
        for child in element.find_all(True):
            if not isinstance(child, Tag) or child.name in SKIP_TAGS:
                continue
            if str(child.get("data-type") or "") == "arch-node":
                self._render_arch_node(box, child)
                nodes += 1
            elif has_class(child, "arch-line-h", "arch-line-v"):
                self._render_arch_line(box, child)
            elif has_class(child, "arch-tip"):
                self._render_arch_tip(box, child)
            elif has_class(child, "arch-edge-label") or str(child.get("data-type") or "") == "arch-edge":
                self._render_arch_label(box, child)
        if not nodes:
            self.ctx.report.warn('Schema sans noeud data-type="arch-node": seul le cadre est exporte.')

    def _node_box(self, container: Box, element: Tag, default_w: float, default_h: float) -> Box:
        """Resolve the box of a positioned child, in the container reference.

        Percentages are relative to the container, as documented for
        ``data-x`` / ``data-y``, never to the slide.
        """
        props = style_props(element)
        x = first_pct(element.get("data-x"), props.get("left"))
        y = first_pct(element.get("data-y"), props.get("top"))
        w = first_pct(element.get("data-width"), props.get("width"), default=default_w)
        h = first_pct(element.get("data-height"), props.get("height"), default=default_h)
        box = container.pct(x, y, w, h)
        return apply_transform(box, props)

    def _render_arch_node(self, container: Box, node: Tag) -> None:
        """Render one diagram node as a filled, outlined and labelled shape."""
        box = self._node_box(container, node, default_w=20.0, default_h=12.0)
        props = style_props(node)
        shape_name = str(node.get("data-shape") or "box").lower()
        shape = self._rect(
            box,
            fill=self._element_color(node) or "FFFFFF",
            shape=_SHAPES.get(shape_name, MSO_SHAPE.ROUNDED_RECTANGLE),
            radius_px=parse_px(props.get("border-radius", "")) or 4.0,
        )
        merged = self.res.props(node)
        line = color_of(merged, ("border", "border-color"), self.res)
        shape.line.color.rgb = RGBColor.from_string(line or self.theme.primary)
        shape.line.width = Pt(1.25)
        frame = shape.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.margin_left = frame.margin_right = Inches(0.04)
        frame.margin_top = frame.margin_bottom = Inches(0.02)
        color = parse_color(merged.get("color"), self.res.css_vars) or self.theme.text
        base = TextStyle(size=11, bold=True, color=color, align="center", space_after=0)
        self._write(frame, self._paragraphs(node, base) or [_Para(base, [(str(node.get("data-label") or ""), base)])])

    def _render_arch_line(self, container: Box, edge: Tag) -> None:
        """Render an edge segment as a thin filled rectangle."""
        props = style_props(edge)
        thickness = 1.5 * PX_IN
        horizontal = has_class(edge, "arch-line-h")
        x = parse_pct(props.get("left", "0%"))
        y = parse_pct(props.get("top", "0%"))
        length = parse_pct(props.get("width" if horizontal else "height", "0%"))
        if length <= 0:
            return
        if horizontal:
            box = container.pct(x, y, length, 0.1)
            box = Box(box.left, box.top - thickness / 2, box.width, thickness)
        else:
            box = container.pct(x, y, 0.1, length)
            box = Box(box.left - thickness / 2, box.top, thickness, box.height)
        self._rect(box, fill=self._edge_color(edge))

    def _render_arch_tip(self, container: Box, tip: Tag) -> None:
        """Render an arrow head as a small rotated triangle."""
        props = style_props(tip)
        size = 8.0 * PX_IN
        x = parse_pct(props.get("left", "0%"))
        y = parse_pct(props.get("top", "0%"))
        anchor = container.pct(x, y, 0.1, 0.1)
        rotations = {"arch-tip-r": 90, "arch-tip-l": 270, "arch-tip-d": 180, "arch-tip-u": 0}
        rotation = next(
            (deg for name, deg in rotations.items() if has_class(tip, name)),
            90,
        )
        box = Box(anchor.left - size / 2, anchor.top - size / 2, size, size)
        shape = self._rect(box, fill=self._edge_color(tip), shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
        shape.rotation = rotation

    def _render_arch_label(self, container: Box, label: Tag) -> None:
        """Render an edge label, or a text based arrow, inside the diagram."""
        text = label.get_text(" ", strip=True)
        if not text:
            return
        props = style_props(label)
        x = parse_pct(props.get("left", "0%"))
        y = parse_pct(props.get("top", "0%"))
        width_pct = parse_pct(props.get("width", "0%")) or 18.0
        box = container.pct(x, y, width_pct, 0.1)
        box = Box(box.left, box.top - 0.09, max(box.width, 0.5), 0.2)
        box = apply_transform(box, props)
        style = self.res.apply(
            TextStyle(size=9.5, color="secondary", space_after=0, align="center"),
            self.res.props(label),
        )
        _, frame = self._textbox(box, anchor="middle")
        self._write(frame, [_Para(style, [(text, style)])])

    def _edge_color(self, edge: Tag) -> str:
        """Resolve the stroke color of an edge element."""
        color = color_of(
            self.res.props(edge),
            ("background", "background-color", "border-top-color", "border-color", "color"),
            self.res,
        )
        return color or self.theme.primary

    def _render_annotated_image(self, box: Box, element: Tag) -> None:
        """Render an annotated image and place its callouts over the picture."""
        img = element.find("img")
        if not isinstance(img, Tag):
            self.ctx.report.warn("Bloc annotated-image sans <img>: annotations seules exportees.")
            image_box = box
        else:
            width = min(box.width, self._max_width(element, box.width))
            centered = "auto" in style_props(element).get("margin", "")
            frame = box.center_horizontally(width) if centered else box.resize(width=width)
            picture = self._add_picture(frame, str(img.get("src") or ""), mode="fit")
            image_box = (
                Box(
                    inches(picture.left),
                    inches(picture.top),
                    inches(picture.width),
                    inches(picture.height),
                )
                if picture
                else frame
            )
        for annotation in element.find_all(attrs={"data-type": "annotation"}):
            if isinstance(annotation, Tag):
                self._render_annotation(image_box, annotation)

    def _render_annotation(self, image_box: Box, annotation: Tag) -> None:
        """Render one callout positioned in percentages of the image box."""
        text = annotation.get_text(" ", strip=True)
        if not text:
            return
        props = style_props(annotation)
        x = first_pct(annotation.get("data-x"), props.get("left"))
        y = first_pct(annotation.get("data-y"), props.get("top"))
        fill, fg = self._annotation_colors(annotation)
        size = 9.5
        width = min(len(text) * size * 0.56 / 72.0 + 0.16, image_box.width * 0.9)
        height = 0.24
        box = Box(
            image_box.left + image_box.width * x / 100.0,
            image_box.top + image_box.height * y / 100.0,
            width,
            height,
        )
        self._rect(box, fill=fill)
        style = TextStyle(size=size, bold=True, color=fg, align="center", space_after=0)
        _, frame = self._textbox(box, anchor="middle")
        self._write(frame, [_Para(style, [(text, style)])])

    def _annotation_colors(self, annotation: Tag) -> tuple[str, str]:
        """Resolve the (fill, text) colors of an annotation."""
        props = self.res.props(annotation)
        fill = color_of(props, FILL_KEYS, self.res)
        color = parse_color(props.get("color"), self.res.css_vars)
        if not fill:
            for name in classes(annotation):
                if name in _ANNOT_CLASS_COLORS:
                    fill, fallback = _ANNOT_CLASS_COLORS[name]
                    return fill, color or fallback
        return fill or "FFF176", color or "262626"

    def _render_image(self, box: Box, element: Tag) -> None:
        """Render a standalone image, fitted inside its block."""
        self._add_picture(box, str(element.get("src") or ""), mode="fit")

    # -- text plumbing -----------------------------------------------------

    def _paragraphs(self, element: Tag, base: TextStyle, level: int = 0) -> list[_Para]:
        """Flatten an element into a list of styled paragraphs."""
        if element.name in SKIP_TAGS:
            return []
        if element.name in {"ul", "ol"}:
            paras: list[_Para] = []
            ordered = element.name == "ol"
            for position, item in enumerate(element.find_all("li", recursive=False), start=1):
                style = self.res.style(item, base)
                item_paras = self._paragraphs(item, style, level + 1)
                if item_paras:
                    marker = f"{position}. " if ordered else "\u2022 "
                    item_paras[0].runs.insert(0, (marker, replace(style, bold=True)))
                paras.extend(item_paras)
            return paras
        if self._has_block_children(element):
            paras = []
            for child in element.children:
                if isinstance(child, NavigableString):
                    text = re.sub(r"\s+", " ", str(child)).strip()
                    if text:
                        paras.append(_Para(base, [(text, base)], level))
                elif isinstance(child, Tag) and child.name not in SKIP_TAGS:
                    style = self.res.style(child, base)
                    if child.name in INLINE_BLOCK_TAGS:
                        style = self._inline_style(child, base)
                    paras.extend(self._paragraphs(child, style, level))
            return paras
        runs = self._runs(element, base)
        return [_Para(base, runs, level)] if runs else []

    def _runs(self, element: Tag, style: TextStyle) -> list[tuple[str, TextStyle]]:
        """Flatten the inline content of an element into styled runs.

        Adjacent inline elements are separated by a space: the templates lay
        them out with a flex ``gap`` that has no equivalent in a PPTX run.
        """
        runs: list[tuple[str, TextStyle]] = []
        previous_was_tag = False
        for node in element.children:
            if isinstance(node, NavigableString):
                text = re.sub(r"[ \t\r\n]+", " ", str(node))
                if text.strip():
                    runs.append((text, style))
                    previous_was_tag = False
            elif isinstance(node, Tag):
                if node.name == "br":
                    runs.append(("\n", style))
                    previous_was_tag = False
                elif node.name in SKIP_TAGS:
                    continue
                else:
                    if previous_was_tag and runs and not runs[-1][0].endswith(" "):
                        runs.append((" ", style))
                    runs.extend(self._runs(node, self._inline_style(node, style)))
                    previous_was_tag = True
        if not runs:
            text = element.get_text(" ", strip=True)
            if text:
                runs.append((text, style))
        return runs

    def _inline_style(self, node: Tag, parent: TextStyle) -> TextStyle:
        """Derive the style of an inline element from its parent style."""
        style = parent
        if node.name in {"strong", "b"}:
            style = replace(style, bold=True)
        if node.name in {"em", "i"}:
            style = replace(style, italic=True)
        if node.name in {"small"} or has_class(node, "arch-node-label", "gantt-sub"):
            style = replace(style, size=max(style.size * 0.78, 7.0), bold=False)
        if node.name in {"code", "kbd"} or has_class(node, "cds-code"):
            style = replace(style, mono=True, size=max(style.size * 0.92, 7.0))
        if node.name == "span" and has_class(node.parent, "slide-h1"):
            style = replace(style, color=self.theme.accent)
        if classes(node):
            style = self.res.style(node, style)
        return self.res.apply(style, self.res.props(node))

    def _write(self, frame: Any, paras: list[_Para]) -> None:
        """Write paragraphs into a text frame, honouring styles and breaks.

        Does not touch ``word_wrap``: the frame's own setting (``_textbox``'s
        default is wrapping on, callers needing a single-line label such as
        a Gantt bar pass ``wrap=False``) must survive writing the runs.
        """
        first = True
        for para in paras:
            chunks = split_runs(para.runs)
            for chunk in chunks:
                paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                first = False
                paragraph.alignment = _ALIGN.get(para.style.align, PP_ALIGN.LEFT)
                paragraph.line_spacing = para.style.line_spacing
                paragraph.space_after = Pt(para.style.space_after)
                if para.level > 1:
                    # Markers are written explicitly, so only nested lists indent.
                    paragraph.level = min(para.level - 1, 4)
                for text, style in chunk:
                    run = paragraph.add_run()
                    run.text = text.upper() if style.upper else text
                    font = run.font
                    font.name = self.theme.mono_font if style.mono else self.theme.font
                    font.size = Pt(style.size)
                    font.bold = style.bold
                    font.italic = style.italic
                    font.color.rgb = RGBColor.from_string(self.theme.color(style.color))

    def _textbox(self, box: Box, anchor: str = "top", wrap: bool = True) -> tuple[Any, Any]:
        """Create a margin free textbox and return ``(shape, text_frame)``.

        ``wrap=False`` also turns off ``spAutoFit``: python-pptx's default
        auto-size only resizes the *shape* to fit unwrapped text (an
        instruction PowerPoint applies live, but LibreOffice's headless
        renderer used for exports and screenshots does not), so a caller
        that pre-fits its text to the box (the Gantt bar labels) needs the
        frame to just clip at the box bounds instead, matching what it
        measured against.
        """
        shape = self.slide.shapes.add_textbox(*box.emu())
        frame = shape.text_frame
        frame.word_wrap = wrap
        if not wrap:
            frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE if anchor == "middle" else MSO_ANCHOR.TOP
        return shape, frame

    # -- shapes and pictures ----------------------------------------------

    def _rect(  # noqa: PLR0913 - one flag per optional shape trait, all keyword-only
        self,
        box: Box,
        fill: str | None = None,
        line: str | None = None,
        *,
        shape: Any = MSO_SHAPE.RECTANGLE,
        radius_px: float | None = None,
        hatch: bool = False,
    ) -> Any:
        """Add an auto shape with a solid fill and no outline by default.

        ``hatch`` swaps the solid fill for a native PPTX diagonal pattern
        fill (foreground ``fill``, white background), the closest built-in
        equivalent of the CSS ``repeating-linear-gradient(45deg, ...)`` hatch
        the templates use to mark a task as done/in a different state: real
        pattern geometry survives the export instead of a flat color or a
        text workaround.
        """
        auto = self.slide.shapes.add_shape(shape, *box.emu())
        drop_theme_style(auto)
        if radius_px and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            smallest = max(min(box.width, box.height) * 72.0, 1.0)
            auto.adjustments[0] = min(max(radius_px / smallest, 0.0), 0.5)
        if hatch and fill:
            auto.fill.patterned()
            auto.fill.pattern = MSO_PATTERN_TYPE.WIDE_UPWARD_DIAGONAL
            auto.fill.fore_color.rgb = RGBColor.from_string(fill)
            auto.fill.back_color.rgb = RGBColor.from_string("FFFFFF")
        elif fill:
            auto.fill.solid()
            auto.fill.fore_color.rgb = RGBColor.from_string(fill)
        else:
            auto.fill.background()
        if line:
            auto.line.color.rgb = RGBColor.from_string(line)
            auto.line.width = Pt(0.75)
        else:
            auto.line.fill.background()
        auto.text_frame.word_wrap = True
        return auto

    def _element_color(self, element: Tag) -> str | None:
        """Resolve the background color of an element (CSS, inline, data-color)."""
        color = color_of(self.res.props(element), FILL_KEYS, self.res)
        return color or parse_color(str(element.get("data-color") or ""), self.res.css_vars)

    def _image_source(self, src: str) -> io.BytesIO | str | None:
        """Resolve an image source to something python-pptx can read.

        Data URIs are decoded in memory; relative paths resolve against the
        directory of the HTML file; remote URLs are refused.
        """
        if not src:
            return None
        if src.startswith("data:"):
            header, _, payload = src.partition(",")
            if "base64" not in header or not payload:
                self.ctx.report.warn("Image data URI non base64 ignoree.")
                return None
            try:
                return io.BytesIO(base64.b64decode(payload, validate=False))
            except (binascii.Error, ValueError):
                self.ctx.report.warn("Image data URI illisible ignoree.")
                return None
        if src.startswith(("http://", "https://", "//")):
            self.ctx.report.warn(f"Image distante non telechargee: {src[:60]}")
            return None
        path = Path(src)
        if not path.is_absolute():
            path = self.ctx.base_dir / path
        if not path.exists():
            self.ctx.report.warn(f"Image introuvable, ignoree: {path}")
            return None
        return str(path)

    def _image_aspect(self, src: str) -> float | None:
        """Return the width / height ratio of an image, ``None`` if unknown."""
        source = self._image_source(src)
        if source is None:
            return None
        try:
            picture = self.slide.shapes.add_picture(source, 0, 0)
        except Exception as exc:
            self.ctx.report.warn(f"Image non decodable ignoree: {exc}")
            return None
        aspect = picture.width / picture.height if picture.height else None
        picture._element.getparent().remove(picture._element)
        return aspect

    def _add_picture(self, box: Box, src: str, mode: str = "fit") -> Any:
        """Insert a picture in ``box``.

        Modes: ``fit`` keeps the aspect ratio inside the box and centers it,
        ``cover`` fills the box and crops the overflow (CSS ``object-fit:
        cover``), ``fit-height`` keeps the height and derives the width,
        ``stretch`` forces the box dimensions.
        """
        source = self._image_source(src)
        if source is None:
            return None
        if isinstance(source, io.BytesIO):
            source.seek(0)
        try:
            picture = self.slide.shapes.add_picture(source, Inches(box.left), Inches(box.top))
        except Exception as exc:
            self.ctx.report.warn(f"Image non inseree ({exc}): {src[:40]}")
            return None
        natural = picture.width / picture.height if picture.height else 1.0

        if mode == "stretch":
            picture.width, picture.height = Inches(box.width), Inches(box.height)
            return picture
        if mode == "fit-height":
            picture.height = Inches(box.height)
            picture.width = Inches(box.height * natural)
            return picture
        if mode == "cover":
            picture.width, picture.height = Inches(box.width), Inches(box.height)
            target = box.width / box.height
            if natural > target:
                crop = (1.0 - target / natural) / 2.0
                picture.crop_left = crop
                picture.crop_right = crop
            elif natural < target:
                crop = (1.0 - natural / target) / 2.0
                picture.crop_top = crop
                picture.crop_bottom = crop
            return picture
        scale = min(box.width / natural, box.height) if natural else box.height
        width, height = scale * natural, scale
        picture.width, picture.height = Inches(width), Inches(height)
        picture.left = Inches(box.left + (box.width - width) / 2.0)
        picture.top = Inches(box.top + (box.height - height) / 2.0)
        return picture


__all__ = ["ExportReport", "find_slides", "to_pptx"]
